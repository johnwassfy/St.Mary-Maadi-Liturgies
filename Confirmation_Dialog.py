from PyQt5.QtWidgets import QDialog, QWidget, QLabel, QCheckBox, QVBoxLayout, QHBoxLayout, QPushButton, QComboBox, QLineEdit, QFrame, QGraphicsDropShadowEffect, QScrollArea, QRadioButton, QButtonGroup, QToolTip
from PyQt5.QtGui import QIcon, QFont, QColor, QCursor
from PyQt5.QtCore import Qt, QSize, QTimer
from pptx import Presentation
from commonFunctions import relative_path
import qtawesome as qta

class SynaxarSection(QWidget):
    def __init__(self):
        super().__init__()
        
        # Section title
        self.title_label = QLabel("السنكسار")
        self.title_label.setStyleSheet("font-weight: bold; color: white; font-size: 16px; margin-bottom: 8px;")
        
        # Radio button group
        self.button_group = QButtonGroup()
        
        # Radio button 1
        self.radio1 = QRadioButton("عناوين فقط")
        self.radio1.setChecked(True)  # Default selection
        self.radio1.setStyleSheet("""
            QRadioButton {
                color: white;
                font-size: 14px;
                padding: 5px;
            }
            QRadioButton::indicator {
                width: 20px;
                height: 20px;
                border: 3px solid white;
                border-radius: 12px;
                background: transparent;
            }
            QRadioButton::indicator:unchecked {
                background: transparent;
                border: 3px solid rgba(255, 255, 255, 180);
            }
            QRadioButton::indicator:checked {
                background: qlineargradient(
                    x1: 0, y1: 0, x2: 1, y2: 1,
                    stop: 0 #00ff88,
                    stop: 1 #00cc66
                );
                border: 3px solid #00ff88;
            }
            QRadioButton::indicator:hover {
                border: 3px solid rgba(255, 255, 255, 255);
            }
        """)
        
        # Radio button 2 (disabled with tooltip)
        self.radio2 = QRadioButton("النص كامل")
        self.radio2.setEnabled(False)  # Disable the radio button
        self.radio2.setToolTip("سيتوفر في التحديثات القادمة")
        self.radio2.setStyleSheet("""
            QRadioButton {
                color: rgba(255, 255, 255, 100);
                font-size: 14px;
                padding: 5px;
            }
            QRadioButton::indicator {
                width: 20px;
                height: 20px;
                border: 3px solid rgba(255, 255, 255, 80);
                border-radius: 12px;
                background: transparent;
            }
            QRadioButton::indicator:unchecked {
                background: transparent;
                border: 3px solid rgba(255, 255, 255, 80);
            }
            QRadioButton::indicator:disabled {
                background: transparent;
                border: 3px solid rgba(255, 255, 255, 50);
            }
            QRadioButton:disabled {
                color: rgba(255, 255, 255, 100);
            }
        """)
        
        # Set custom tooltip style
        self.radio2.setStyleSheet(self.radio2.styleSheet() + """
            QToolTip {
                background-color: rgba(0, 0, 0, 200);
                color: white;
                border: 1px solid rgba(255, 255, 255, 100);
                border-radius: 5px;
                padding: 5px;
                font-size: 12px;
            }
        """)
        
        # Add radio buttons to group
        self.button_group.addButton(self.radio1)
        self.button_group.addButton(self.radio2)
        
        # Connect radio button changes to signal
        self.radio1.toggled.connect(self.on_selection_changed)
        self.radio2.toggled.connect(self.on_selection_changed)
        
        # Layout
        layout = QVBoxLayout()
        layout.setSpacing(8)
        layout.addWidget(self.title_label)
        
        # Radio buttons layout
        radio_layout = QHBoxLayout()
        radio_layout.addWidget(self.radio1)
        radio_layout.addWidget(self.radio2)
        radio_layout.addStretch()
        layout.addLayout(radio_layout)
        
        self.setLayout(layout)
    
    def get_selected_option(self):
        """Returns the selected synaxar option"""
        if self.radio1.isChecked():
            return 1
        elif self.radio2.isChecked():
            return 2
        else:
            return 1  # Default fallback
    
    def on_selection_changed(self):
        """Called when radio button selection changes"""
        # This will be connected to parent dialog's method
        pass

class CustomRow(QWidget):
    def __init__(self, label_text):
        super().__init__()
        
        self.label = QLabel(label_text)
        self.label.setStyleSheet("font-weight: bold; color: white; font-size: 14px;")
        
        self.textbox = QLabel()
        self.combo_box = QComboBox()
        self.combo_box.addItems(["الاسقف", "المطران"])
        self.combo_box.setStyleSheet("""
            QComboBox {
                font-size: 14px;
                border: 2px solid white;
                border-radius: 5px;
                padding: 5px 25px 5px 5px;
                background-color: white;
                color: #1a365d;
                min-width: 80px;
            }
            QComboBox::drop-down {
                subcontrol-origin: padding;
                subcontrol-position: top right;
                width: 20px;
                border-left: 1px solid #1a365d;
                background: white;
                border-top-right-radius: 3px;
                border-bottom-right-radius: 3px;
            }
            QComboBox::down-arrow {
                image: none;
                border-left: 4px solid transparent;
                border-right: 4px solid transparent;
                border-top: 6px solid #1a365d;
                width: 0px;
                height: 0px;
            }
            QComboBox::down-arrow:hover {
                border-top: 6px solid #2a466d;
            }
        """)
        
        self.line_edit = QLineEdit()
        self.line_edit.setPlaceholderText("أدخل الاسم هنا")
        self.line_edit.setStyleSheet("""
            QLineEdit {
                border: 2px solid white;
                border-radius: 5px;
                padding: 5px;
                background-color: white;
                color: #1a365d;
                font-size: 14px;
            }
            QLineEdit:focus {
                border: 2px solid #ffffff;
                background-color: rgba(255, 255, 255, 230);
            }
        """)

        layout = QHBoxLayout()
        layout.addWidget(self.label)        
        layout.addWidget(self.line_edit)
        layout.addWidget(self.combo_box)
        layout.addStretch()  # Add stretch at the end to push everything right
        

        self.setLayout(layout)

class ConfirmationDialog(QDialog):
    def __init__(self, parent=None, coptic_date="", type=""):
        super().__init__(parent)
        
        self.coptic_date = coptic_date
        self.type = type
        dialog_title = self.type + coptic_date if coptic_date else "صلاة"
        self.setWindowTitle(dialog_title)
        self.setFixedSize(450, 300)
        
        # Make dialog modal - will be attached to parent window
        self.setWindowFlags(Qt.Dialog | Qt.FramelessWindowHint | Qt.WindowSystemMenuHint | Qt.WindowTitleHint)
        self.setModal(True)
        
        # Main layout
        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(0, 0, 0, 0)
        main_layout.setSpacing(0)
        
        # Set gradient background for the entire dialog
        self.setStyleSheet("""
            QDialog {
                background: qlineargradient(
                    x1: 0, y1: 0, x2: 1, y2: 1,
                    stop: 0 rgba(26, 54, 93, 245),
                    stop: 0.6 rgba(42, 70, 109, 245),
                    stop: 1 rgba(60, 90, 130, 245)
                );
                border-radius: 10px;
                border: 1px solid rgba(200, 200, 200, 150);
            }
            QLabel {
                font-size: 14px;
                color: white;
            }
        """)
        
        # Header
        header = self.create_header()
        main_layout.addWidget(header)
        
        # Create scroll area for content
        scroll_area = QScrollArea()
        scroll_area.setWidgetResizable(True)
        scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        scroll_area.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        scroll_area.setStyleSheet("""
            QScrollArea {
                background: transparent;
                border: none;
            }
            QScrollBar:vertical {
                background: rgba(255, 255, 255, 100);
                width: 8px;
                border-radius: 4px;
                margin: 0;
            }
            QScrollBar::handle:vertical {
                background: rgba(255, 255, 255, 200);
                border-radius: 4px;
                min-height: 20px;
            }
            QScrollBar::handle:vertical:hover {
                background: rgba(255, 255, 255, 255);
            }
            QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {
                height: 0px;
            }
            QScrollBar::add-page:vertical, QScrollBar::sub-page:vertical {
                background: transparent;
            }
        """)
        
        # Main content widget
        content_widget = QWidget()
        content_widget.setStyleSheet("background: transparent; border: none;")
        content_layout = QVBoxLayout(content_widget)
        content_layout.setContentsMargins(15, 15, 15, 15)
        content_layout.setSpacing(12)

        # Row 1
        self.label1 = QLabel("حضور اسقف الابراشية ")
        self.label1.setStyleSheet("font-weight: bold; font-size: 14px; color: white;")
        self.checkbox1 = QCheckBox()
        self.checkbox1.setStyleSheet("""
            QCheckBox {
                color: white;
                margin-right: 8px;
                font-size: 16px;
            }
            QCheckBox::indicator {
                width: 22px;
                height: 22px;
                border: 3px solid white;
                border-radius: 6px;
                background: transparent;
            }
            QCheckBox::indicator:unchecked {
                background: transparent;
                border: 3px solid rgba(255, 255, 255, 180);
            }
            QCheckBox::indicator:checked {
                background: qlineargradient(
                    x1: 0, y1: 0, x2: 1, y2: 1,
                    stop: 0 #00ff88,
                    stop: 1 #00cc66
                );
                border: 3px solid #00ff88;
            }
            QCheckBox::indicator:hover {
                border: 3px solid rgba(255, 255, 255, 255);
            }
        """)

        # Row 2 & 3
        self.row2 = CustomRow("حضور اسقف ضيف واحد")
        self.row3 = CustomRow("حضور اسقف ضيف ثاني")
        
        # Initially disable row3
        self.set_row3_enabled(False)
        
        # Connect row2 text change to enable/disable row3
        self.row2.line_edit.textChanged.connect(self.on_row2_text_changed)

        # Synaxar section (only for "قداس" type)
        self.synaxar_section = None
        if self.type == "قداس":
            self.synaxar_section = SynaxarSection()

        # Save Button
        self.update_button = QPushButton("تحميل الملف")
        self.update_button.clicked.connect(self.update_powerpoint)
        self.update_button.setStyleSheet("""
            QPushButton {
                background-color: #ffffff;
                color: #1a365d;
                border-radius: 8px;
                padding: 8px 16px;
                font-weight: bold;
                font-size: 14px;
                border: 2px solid #ffffff;
            }
            QPushButton:hover {
                background-color: rgba(255, 255, 255, 200);
                color: #1a365d;
            }
            QPushButton:pressed {
                background-color: rgba(255, 255, 255, 150);
            }
        """)
        shadow = QGraphicsDropShadowEffect()
        shadow.setBlurRadius(12)
        shadow.setOffset(2, 2)
        shadow.setColor(QColor(0, 0, 0, 80))
        self.update_button.setGraphicsEffect(shadow)

        # Row 1 layout
        row1_layout = QHBoxLayout()
        row1_layout.addWidget(self.checkbox1)
        row1_layout.addWidget(self.label1)
        row1_layout.addStretch()  # Add stretch at the end to push everything right
        content_layout.addLayout(row1_layout)

        # Add remaining widgets
        content_layout.addWidget(self.row2)
        content_layout.addWidget(self.row3)
        
        # Add synaxar section if type is "قداس"
        if self.synaxar_section:
            # Add separator line
            separator = QFrame()
            separator.setFrameShape(QFrame.HLine)
            separator.setStyleSheet("color: rgba(255, 255, 255, 100); margin: 10px 0;")
            content_layout.addWidget(separator)
            content_layout.addWidget(self.synaxar_section)

        # Button layout
        button_layout = QHBoxLayout()
        button_layout.addWidget(self.update_button)
        button_layout.addStretch()  # Add stretch at the end to push button right
        content_layout.addLayout(button_layout)

        # Set the content widget to scroll area and add to main layout
        scroll_area.setWidget(content_widget)
        main_layout.addWidget(scroll_area, 1)
        
        # Set Arabic RTL layout
        self.setLayoutDirection(Qt.RightToLeft)

    def on_row2_text_changed(self, text):
        """Enable/disable row3 based on row2 text content"""
        has_text = bool(text.strip())
        self.set_row3_enabled(has_text)

    def set_row3_enabled(self, enabled):
        """Enable or disable row3 with visual feedback"""
        self.row3.line_edit.setEnabled(enabled)
        self.row3.combo_box.setEnabled(enabled)
        
        if enabled:
            # Active state - normal colors
            self.row3.label.setStyleSheet("font-weight: bold; color: white; font-size: 14px;")
            self.row3.line_edit.setStyleSheet("""
                QLineEdit {
                    border: 2px solid white;
                    border-radius: 5px;
                    padding: 5px;
                    background-color: white;
                    color: #1a365d;
                    font-size: 14px;
                }
                QLineEdit:focus {
                    border: 2px solid #ffffff;
                    background-color: rgba(255, 255, 255, 230);
                }
            """)
            self.row3.combo_box.setStyleSheet("""
                QComboBox {
                    font-size: 14px;
                    border: 2px solid white;
                    border-radius: 5px;
                    padding: 5px 25px 5px 5px;
                    background-color: white;
                    color: #1a365d;
                    min-width: 80px;
                }
                QComboBox::drop-down {
                    subcontrol-origin: padding;
                    subcontrol-position: top right;
                    width: 20px;
                    border-left: 1px solid #1a365d;
                    background: white;
                    border-top-right-radius: 3px;
                    border-bottom-right-radius: 3px;
                }
                QComboBox::down-arrow {
                    image: none;
                    border-left: 4px solid transparent;
                    border-right: 4px solid transparent;
                    border-top: 6px solid #1a365d;
                    width: 0px;
                    height: 0px;
                }
                QComboBox::down-arrow:hover {
                    border-top: 6px solid #2a466d;
                }
            """)
        else:
            # Inactive state - grayed out
            self.row3.label.setStyleSheet("font-weight: bold; color: rgba(255, 255, 255, 100); font-size: 14px;")
            self.row3.line_edit.setStyleSheet("""
                QLineEdit {
                    border: 2px solid rgba(255, 255, 255, 80);
                    border-radius: 5px;
                    padding: 5px;
                    background-color: rgba(255, 255, 255, 50);
                    color: rgba(26, 54, 93, 120);
                    font-size: 14px;
                }
            """)
            self.row3.combo_box.setStyleSheet("""
                QComboBox {
                    font-size: 14px;
                    border: 2px solid rgba(255, 255, 255, 80);
                    border-radius: 5px;
                    padding: 5px 25px 5px 5px;
                    background-color: rgba(255, 255, 255, 50);
                    color: rgba(26, 54, 93, 120);
                    min-width: 80px;
                }
                QComboBox::drop-down {
                    subcontrol-origin: padding;
                    subcontrol-position: top right;
                    width: 20px;
                    border-left: 1px solid rgba(26, 54, 93, 80);
                    background: rgba(255, 255, 255, 50);
                    border-top-right-radius: 3px;
                    border-bottom-right-radius: 3px;
                }
                QComboBox::down-arrow {
                    image: none;
                    border-left: 4px solid transparent;
                    border-right: 4px solid transparent;
                    border-top: 6px solid rgba(26, 54, 93, 120);
                    width: 0px;
                    height: 0px;
                }
            """)
            # Clear the text when disabling
            self.row3.line_edit.clear()

    def create_header(self):
        header = QFrame()
        header.setFixedHeight(50)
        header.setStyleSheet("""
            QFrame {
                background: qlineargradient(
                    x1: 0, y1: 0, x2: 1, y2: 0,
                    stop: 0 #1a365d,
                    stop: 1 #2a466d
                );
                border-top-left-radius: 10px;
                border-top-right-radius: 10px;
            }
        """)
        
        header_layout = QHBoxLayout(header)
        header_layout.setContentsMargins(15, 0, 15, 0)
        
        # Title with icon
        title_layout = QHBoxLayout()
        
        # Title
        if self.coptic_date:
            title_text = f"{self.type} - {self.coptic_date}"
        else:
            title_text = self.type
        title_label = QLabel(title_text)
        title_font = QFont()
        title_font.setPointSize(16)
        title_font.setBold(True)
        title_label.setFont(title_font)
        title_label.setStyleSheet("color: white; background: transparent;")
        title_layout.addWidget(title_label)
        
        # Close button in header
        close_button = QPushButton()
        close_button.setFixedSize(30, 30)
        close_button.setStyleSheet("""
            QPushButton {
                background-color: transparent;
                border: none;
            }
            QPushButton:hover {
                background-color: rgba(255, 0, 0, 150);
                border-radius: 15px;
            }
        """)
        close_button.setCursor(Qt.PointingHandCursor)
        
        # Set X icon
        try:
            close_button.setIcon(qta.icon("fa5s.times", color="white"))
            close_button.setIconSize(QSize(16, 16))
        except:
            close_button.setText("×")
            close_button.setStyleSheet("""
                QPushButton {
                    color: white;
                    font-size: 16pt;
                    font-weight: bold;
                    background-color: transparent;
                    border: none;
                }
                QPushButton:hover {
                    background-color: rgba(255, 0, 0, 150);
                    border-radius: 15px;
                }
            """)
        
        close_button.clicked.connect(self.reject)
        
        header_layout.addLayout(title_layout)
        header_layout.addStretch()
        header_layout.addWidget(close_button)
        
        return header

    def update_powerpoint(self):
        presentation = Presentation(relative_path(r"Data\CopyData\في حضور الاسقف و اساقفة ضيوف.pptx"))

        # Get synaxar option if available
        synaxar_option = None
        if self.synaxar_section:
            synaxar_option = self.synaxar_section.get_selected_option()

        for slide in presentation.slides:
            slide_contains_word = False
            slide_contains_word2 = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "(الضيف)" in run.text:
                                slide_contains_word = True
                            if "(الضيف2)" in run.text:
                                slide_contains_word2 = True

            if slide_contains_word:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "(الضيف)" in run.text and self.row2.line_edit.text():
                                    run.text = run.text.replace("(الضيف)", self.row2.line_edit.text())
                                if self.row2.combo_box.currentText() == "المطران":
                                    run.text = run.text.replace("الأسقف", "المطران")
                                    run.text = run.text.replace("اسقفنا", "مطراننا")
                                    run.text = run.text.replace("ايبيسكوبوس", "متروبوليتيس")
                                    run.text = run.text.replace("إيبيسكوبوتيس", "متروبوليتيس")
                                    run.text = run.text.replace("n`epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("`epickopou tyc", "mmytropolityc")

            if slide_contains_word2:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "(الضيف2)" in run.text and self.row3.line_edit.text():
                                    run.text = run.text.replace("(الضيف2)", self.row3.line_edit.text())
                                if self.row3.combo_box.currentText() == "المطران":
                                    run.text = run.text.replace("الأسقف", "المطران")
                                    run.text = run.text.replace("اسقفنا", "مطراننا")
                                    run.text = run.text.replace("ايبيسكوبوس", "متروبوليتيس")
                                    run.text = run.text.replace("إيبيسكوبوتيس", "متروبوليتيس")
                                    run.text = run.text.replace("n`epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("`epickopou tyc", "mmytropolityc")

        # Handle synaxar option if available
        if synaxar_option:
            # You can add synaxar-specific processing here
            # For example, modify specific slides or text based on the selected option
            pass

        presentation.save(relative_path(r"Data\حضور الأسقف.pptx"))
        # Let the main application handle closing the dialog

    def close_dialog(self):
        """Method to properly close the dialog from external code"""
        self.accept()

    def mousePressEvent(self, event):
        # Allow dragging the frameless window from the header area
        if event.button() == Qt.LeftButton and event.y() < 50:  # 50 is header height
            self._drag_pos = event.globalPos() - self.frameGeometry().topLeft()
            event.accept()

    def mouseMoveEvent(self, event):
        # Move the window with mouse
        if event.buttons() == Qt.LeftButton and hasattr(self, '_drag_pos'):
            self.move(event.globalPos() - self._drag_pos)
            event.accept()


# Keep the old class name as an alias for backward compatibility
Confirm = ConfirmationDialog
