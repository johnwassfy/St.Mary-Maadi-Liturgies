from PyQt5.QtWidgets import QWidget, QLabel, QCheckBox, QVBoxLayout, QHBoxLayout, QPushButton, QComboBox, QLineEdit, QFrame, QGraphicsDropShadowEffect
from PyQt5.QtGui import QIcon, QFont, QColor
from PyQt5.QtCore import Qt
from pptx import Presentation
from commonFunctions import relative_path

class CustomRow(QWidget):
    def __init__(self, label_text):
        super().__init__()
        
        self.label = QLabel(label_text)
        self.label.setStyleSheet("font-weight: bold; color: #1a365d; font-size: 14px;")
        
        self.textbox = QLabel()
        self.combo_box = QComboBox()
        self.combo_box.addItems(["الاسقف", "المطران"])
        self.combo_box.setStyleSheet("""
            QComboBox {
                font-size: 14px;
                border: 1px solid #1a365d;
                border-radius: 5px;
                padding: 5px;
                background-color: white;
                min-width: 80px;
            }
        """)
        
        self.line_edit = QLineEdit()
        self.line_edit.setPlaceholderText("أدخل الاسم هنا")
        self.line_edit.setStyleSheet("""
            QLineEdit {
                border: 1px solid #1a365d;
                border-radius: 5px;
                padding: 5px;
                background-color: white;
                font-size: 14px;
            }
        """)

        layout = QHBoxLayout()
        layout.addStretch()
        layout.addWidget(self.combo_box)
        layout.addWidget(self.line_edit)
        layout.addWidget(self.label)

        self.setLayout(layout)

class Bishop(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("حضور الأسقف")
        self.setGeometry(550, 300, 400, 220)
        self.setWindowIcon(QIcon(relative_path(r"Data\الصور\Logo.ico")))
        self.setStyleSheet("""
            QWidget {
                background-color: #f0f5ff;
                font-family: 'Segoe UI';
            }
            QLabel {
                font-size: 14px;
                color: #1a365d;
            }
        """)

        # Row 1
        self.label1 = QLabel("حضور اسقف الابراشية ")
        self.label1.setStyleSheet("font-weight: bold; font-size: 14px;")
        self.checkbox1 = QCheckBox()
        self.checkbox1.setStyleSheet("margin-right: 8px;")

        # Row 2 & 3
        self.row2 = CustomRow("حضور اسقف ضيف واحد")
        self.row3 = CustomRow("حضور اسقف ضيف ثاني")

        # Save Button
        self.update_button = QPushButton("حفـظ")
        self.update_button.clicked.connect(self.update_powerpoint)
        self.update_button.setStyleSheet("""
            QPushButton {
                background-color: #1a365d;
                color: white;
                border-radius: 8px;
                padding: 8px 16px;
                font-weight: bold;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: #2a466d;
            }
            QPushButton:pressed {
                background-color: #0a264d;
            }
        """)
        shadow = QGraphicsDropShadowEffect()
        shadow.setBlurRadius(12)
        shadow.setOffset(2, 2)
        shadow.setColor(QColor(0, 0, 0, 80))
        self.update_button.setGraphicsEffect(shadow)

        # Layouts
        layout = QVBoxLayout()
        layout.setSpacing(12)
        layout.setContentsMargins(15, 15, 15, 15)

        # Row 1 layout
        row1_layout = QHBoxLayout()
        row1_layout.addStretch()
        row1_layout.addWidget(self.checkbox1)
        row1_layout.addWidget(self.label1)
        row1_layout.addStretch()
        layout.addLayout(row1_layout)

        # Add remaining widgets
        layout.addWidget(self.row2)
        layout.addWidget(self.row3)

        # Button layout
        button_layout = QHBoxLayout()
        button_layout.addStretch()
        button_layout.addWidget(self.update_button)
        layout.addLayout(button_layout)

        self.setLayout(layout)

    def update_powerpoint(self):
        presentation = Presentation(relative_path(r"Data\CopyData\في حضور الاسقف و اساقفة ضيوف.pptx"))

        for slide in presentation.slides:
            slide_contains_word = False
            slide_contains_word2 = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "(الضيف)" in run.text:
                                slide_contains_word = True
                            if "(الضيف2)" in run.text:
                                slide_contains_word2 = True

            if slide_contains_word:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "(الضيف)" in run.text and self.row2.line_edit.text():
                                    run.text = run.text.replace("(الضيف)", self.row2.line_edit.text())
                                if self.row2.combo_box.currentText() == "المطران":
                                    run.text = run.text.replace("الأسقف", "المطران")
                                    run.text = run.text.replace("اسقفنا", "مطراننا")
                                    run.text = run.text.replace("ايبيسكوبوس", "متروبوليتيس")
                                    run.text = run.text.replace("إيبيسكوبوتيس", "متروبوليتيس")
                                    run.text = run.text.replace("n`epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("`epickopou tyc", "mmytropolityc")

            if slide_contains_word2:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "(الضيف2)" in run.text and self.row3.line_edit.text():
                                    run.text = run.text.replace("(الضيف2)", self.row3.line_edit.text())
                                if self.row3.combo_box.currentText() == "المطران":
                                    run.text = run.text.replace("الأسقف", "المطران")
                                    run.text = run.text.replace("اسقفنا", "مطراننا")
                                    run.text = run.text.replace("ايبيسكوبوس", "متروبوليتيس")
                                    run.text = run.text.replace("إيبيسكوبوتيس", "متروبوليتيس")
                                    run.text = run.text.replace("n`epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("epickopoc", "mmytropolityc")
                                    run.text = run.text.replace("`epickopou tyc", "mmytropolityc")

        presentation.save(relative_path(r"Data\حضور الأسقف.pptx"))
