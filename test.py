# import win32com.client

# def print_custom_layouts(master_slide):
#     """
#     Print information about the custom layouts in the given master slide.

#     Args:
#     - master_slide: PowerPoint master slide object.
#     """
#     for i, layout in enumerate(master_slide.CustomLayouts, start=1):
#         print(f"Layout {i}:")
#         print(f"Name: {layout.Name}")
#         print(f"Index: {i}")
#         # You can print and inspect other properties of the layout if needed

# def main():
#     # Initialize PowerPoint application
#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#     powerpoint.Visible = True  # Open PowerPoint application

#     # Open the PowerPoint presentation containing the master slide
#     presentation_path = r"C:\Users\dell\Desktop\5dmt Shashat\Codes and Files\قداس.pptx"  # Change this to your presentation path
#     presentation = powerpoint.Presentations.Open(presentation_path)

#     # Get the master slide
#     master_slide = presentation.Designs(1).SlideMaster

#     # Print information about the custom layouts
#     print_custom_layouts(master_slide)

#     # Close the presentation and PowerPoint application
#     presentation.Close()
#     powerpoint.Quit()

# if __name__ == "__main__":
#     main()

# def find_master_slide_by_name(presentation, master_name):
#     for design in presentation.Designs:
#         for master in design.SlideMaster.CustomLayouts:
#             if master.Name == master_name:
#                 return master
#     return None

# # Open the PowerPoint presentation
# powerpoint = win32com.client.Dispatch("PowerPoint.Application")
# presentation_path = r"C:\Users\dell\Desktop\5dmt Shashat\Codes and Files\قداس.pptx"
# presentation = powerpoint.Presentations.Open(presentation_path)

# master_name = "Blank3"  # Change this to the name of the master slide you want to find
# default_master = find_master_slide_by_name(presentation, master_name)
# if default_master:
#     print("Found the master slide:", default_master.Name)
# else:
#     print("Master slide not found.")


# from pptx import Presentation

# def find_slide_index_by_text(presentation_path, text_to_find, start_index=1):
#     prs = Presentation(presentation_path)
#     num_slides = len(prs.slides)
#     if start_index < 1 or start_index > num_slides:
#         raise ValueError("start_index is out of range")

#     for i in range(start_index - 1, num_slides):
#         slide = prs.slides[i]
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and text_to_find in shape.text:
#                 return i + 1
#     return -1
# # Example usage:
# presentation_path = r"قداس.pptx"
# text_to_find = "إبن الله دخل أورشليم."
# slide_index = find_slide_index_by_text(presentation_path, text_to_find, 3469)
# print("Slide index:", slide_index)

# import win32com.client

# def find_slide_index_by_text(presentation_path, text_to_find, start_index=1):
#     ppt_app = win32com.client.Dispatch("PowerPoint.Application")
#     ppt_app.Visible = True
#     presentation = ppt_app.Presentations.Open(presentation_path)
#     num_slides = presentation.Slides.Count
#     if start_index < 1 or start_index > num_slides:
#         raise ValueError("start_index is out of range")

#     for i in range(start_index, num_slides + 1):
#         slide = presentation.Slides(i)
#         for shape in slide.Shapes:
#             if shape.HasTextFrame:
#                 text_frame = shape.TextFrame
#                 if text_frame.HasText:
#                     if text_to_find in text_frame.TextRange.Text:
#                         return i
#     return -1

# # Example usage:
# presentation_path = r"C:\Users\dell\Desktop\5dmt Shashat\Codes and Files\قداس.pptx"
# text_to_find = "إبن الله دخل أورشليم."
# slide_index = find_slide_index_by_text(presentation_path, text_to_find, 3469)
# print("Slide index:", slide_index)

# from pptx import Presentation
# import win32com.client
# import time

# start = time.time()

# # Get the layouts and slide indices for the master slide

# def get_layouts_and_slide_indices(presentation_path, ms):
#     prs = Presentation(presentation_path)
#     master_slide = prs.slide_masters[ms]  # Index 1 corresponds to the second master slide

#     layouts_and_indices = {}
    
#     # Find the layouts in the master slide
#     for layout in master_slide.slide_layouts:
#         layout_name = layout.name
#         layout_indices = []
        
#         # Find slides using the current layout
#         for i, slide in enumerate(prs.slides, start=1):
#             if slide.slide_layout == layout:
#                 layout_indices.append(i)
        
#         # Store layout name and indices
#         layouts_and_indices[layout_name] = layout_indices

#     return layouts_and_indices

# def update_layout_for_slides(presentation_path, ms_source, ms_target):
#     # Create PowerPoint application object
#     ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    
#     # Open the presentation
#     presentation = ppt_app.Presentations.Open(presentation_path)
    
#     # Get the target master
#     target_master = presentation.Designs(ms_target).SlideMaster
    
#     # Get layouts and slide indices from MS2
#     layouts_and_indices_ms2 = get_layouts_and_slide_indices(presentation_path, ms_source)
    
#     # Iterate through the layouts in the target master (MS4)
#     for target_layout in target_master.CustomLayouts:
#         # Get the layout name
#         layout_name = target_layout.Name
        
#         # Check if the layout exists in the MS2 layouts
#         if layout_name in layouts_and_indices_ms2:
#             # Get the indices of slides using this layout in MS2
#             slide_indices_to_change = layouts_and_indices_ms2[layout_name]
            
#             # Iterate through the slides and switch the layout
#             for slide_index in slide_indices_to_change:
#                 slide = presentation.Slides(slide_index)
#                 slide.CustomLayout = target_layout

# presentation_path = r"C:\Users\dell\Desktop\5dmt Shashat\Codes and Files\قداس.pptx"

# # Update layouts for slides using master slide 2 to match those in master slide 4
# update_layout_for_slides(presentation_path, 1, 4)

# end = time.time()

# print(end - start)

# import win32com.client

# def inches_to_points(inches):
#     return inches * 72

# def add_picture_to_slides(presentation_path, picture_path):
#     # Create an instance of PowerPoint
#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    
#     # Show PowerPoint application
#     powerpoint.Visible = True
    
#     # Open an existing presentation
#     presentation = powerpoint.Presentations.Open(presentation_path)
    
#     # Iterate over slides and add picture to each slide
#     for slide in presentation.Slides:
#         # Insert picture into the slide
#         pic_shape = slide.Shapes.AddPicture(picture_path, LinkToFile=False, SaveWithDocument=True, Left=0, Top=144, Width=720, Height=396)
#         pic_shape.ZOrder(1)
    
#     # Save the presentation
#     presentation.Save()
    
#     # Close the presentation and PowerPoint application
#     presentation.Close()
#     powerpoint.Quit()

# from pptx import Presentation

# def remove_background_picture(pptx_file):
#     prs = Presentation(pptx_file)
    
#     # Check if the first slide has a picture in it
#     first_slide = prs.slides[0]
#     shapes = first_slide.shapes
#     has_background_picture = False
    
#     # Check if there's a picture in the background
#     for shape in shapes:
#         if shape.shape_type == 13:  # Shape type 13 represents a picture
#             has_background_picture = True
#             break
    
#     # If the first slide has a picture, proceed to remove background pictures from all slides
#     if has_background_picture:
#         for slide in prs.slides:
#             shapes = slide.shapes
#             background_picture = None
            
#             # Find the picture that's in the background
#             for shape in shapes:
#                 if shape.shape_type == 13:  # Shape type 13 represents a picture
#                     if background_picture is None or shape.left < background_picture.left:
#                         background_picture = shape
            
#             # Remove the background picture
#             if background_picture is not None:
#                 background_picture._element.getparent().remove(background_picture._element)    
#     else:
#         print("The first slide does not have a background picture. Exiting without removing background pictures from other slides.")
    
#     prs.save(pptx_file)  

# # Example usage:
# presentation_path = r"C:\Users\dell\Desktop\5dmt Shashat\Codes and Files\باكر.pptx"
# # picture_path = r"C:\Users\dell\Desktop\5dmt Shashat\Codes and Files\Data\Designs\الشعانين.png"
# # # add_picture_to_slides(presentation_path, picture_path)
# remove_background_picture(presentation_path)

# from pptx import Presentation

# def print_slide_words(presentation_path, slide_number):
#     # Load the PowerPoint presentation
#     prs = Presentation(presentation_path)

#     # Check if the slide number is valid
#     if slide_number < 0 or slide_number >= len(prs.slides):
#         print(f"Slide number {slide_number} is out of range.")
#         return

#     # Get the specified slide
#     slide = prs.slides[slide_number]

#     # Initialize a list to store words
#     words = []

#     # Iterate through shapes in the slide
#     for shape in slide.shapes:
#         if not shape.has_text_frame:
#             continue
#         text_frame = shape.text_frame
#         for paragraph in text_frame.paragraphs:
#             for run in paragraph.runs:
#                 # Split text into words and add to list
#                 words.extend(run.text.split())

#     # Print each word in a new line
#     for word in words:
#         print(word)

# # Example usage:
# presentation_path = 'قداس.pptx'  # Replace with your presentation file path
# slide_number = 120  # Replace with the slide number you want to extract words from

# print_slide_words(presentation_path, slide_number)

# import win32com.client

# def open_or_goto_slide(ppt_path, slide_index, slideshow):
#     # Initialize PowerPoint application
#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    
#     # Set up a flag to check if the presentation is already open
#     is_open = False
#     presentation = None

#     # Iterate over the currently open presentations
#     for pres in powerpoint.Presentations:
#         if pres.FullName == ppt_path:  # Check if the file path matches
#             is_open = True
#             presentation = pres
#             break
    
#     if not is_open:
#         presentation = powerpoint.Presentations.Open(ppt_path)
    
#     # Check if the presentation is already in slideshow mode
#     is_showing = False
#     for slide_show in powerpoint.SlideShowWindows:
#         if slide_show.Presentation == presentation:
#             is_showing = True
#             break

#     if is_showing and slideshow:
#         # If in slideshow mode and slideshow is true, just go to the specified slide
#         slide_show.View.GotoSlide(slide_index)
#     elif not is_showing and not slideshow:
#         # If not in slideshow mode, select the slide
#         slide = presentation.Slides(slide_index)
#         slide.Select()
#     elif not is_showing and slideshow:
#         # Start slideshow if not currently in slideshow mode
#         slide_show = presentation.SlideShowSettings.Run()
#         slide_show.View.GotoSlide(slide_index)

#     # Make sure PowerPoint window is visible
#     powerpoint.Visible = True

# # Provide the full path to your PowerPoint file and the slide index
# ppt_path = r"C:\Users\dell\Desktop\5dmt Shashat\Codes and Files\Data\كتاب المدائح.pptx"
# slide_index = 2  # Go to the 3rd slide
# open_or_goto_slide(ppt_path, slide_index, True)

# from pptx import Presentation
# def hide_slide(ppt_file: str, slide_number: int):
#     # Open the presentation
#     presentation = Presentation(ppt_file)
    
#     # Adjust for 0-based indexing of slides
#     slide_index = slide_number - 1
    
#     # Check if the slide number is valid
#     if 0 <= slide_index < len(presentation.slides):
#         # Hide the slide
#         slide = presentation.slides[slide_index]
#         slide._element.set('show', '1')  # Sets the slide to be hidden
        
#         # Save the modified presentation
#         presentation.save(ppt_file)
#         print(f"Slide {slide_number} has been hidden and saved to {ppt_file}.")
#     else:
#         print(f"Invalid slide number: {slide_number}.")

# from pptx import Presentation
# def get_slide_ids_by_number_pairs(pptx_path, slide_number_pairs):
#     # Load the presentation
#     prs = Presentation(pptx_path)
#     num_slides = len(prs.slides)
    
#     # Initialize the list to hold pairs of SlideIDs
#     slide_id_pairs = []

#     for slide_num1, slide_num2 in slide_number_pairs:
#         # Check if both slide numbers are within the valid range
#         if slide_num1 < 1 or slide_num1 > num_slides or slide_num2 < 1 or slide_num2 > num_slides:
#             print(f"Error: Slide numbers ({slide_num1}, {slide_num2}) out of range.")
#             slide_id_pairs.append((None, None))  # Append None for invalid pairs
#             continue
        
#         # Get the SlideID for each slide number
#         slide_id1 = prs.slides[slide_num1 - 1].slide_id  # Convert to 0-based index
#         slide_id2 = prs.slides[slide_num2 - 1].slide_id

#         # Add the SlideID pair to the list
#         slide_id_pairs.append((slide_id1, slide_id2))
    
#     return slide_id_pairs

import pptx
import pandas as pd

def extract_slide_layouts_to_excel(pptx_file, output_excel_file):
    # Load the PowerPoint presentation
    prs = pptx.Presentation(pptx_file)
    
    # Data structure to store layout names and slide ranges
    layout_data = []

    # Access the single slide master and iterate over its layouts
    slide_master = prs.slide_master
    for layout in slide_master.slide_layouts:
        layout_name = layout.name if layout.name else "Unnamed Layout"
        slide_numbers = []
        
        # Collect slide numbers for this layout
        for slide_index, slide in enumerate(prs.slides, start=1):
            if slide.slide_layout == layout:
                slide_numbers.append(slide_index)
        
        # Convert slide numbers into ranges
        slide_ranges = []
        if slide_numbers:
            start = slide_numbers[0]
            for i in range(1, len(slide_numbers)):
                if slide_numbers[i] != slide_numbers[i - 1] + 1:
                    # Found the end of a range
                    slide_ranges.append(f"{start}-{slide_numbers[i - 1]}" if start != slide_numbers[i - 1] else str(start))
                    start = slide_numbers[i]
            # Add the last range
            slide_ranges.append(f"{start}-{slide_numbers[-1]}" if start != slide_numbers[-1] else str(start))
        
        # Add layout name and slide ranges to the data
        layout_data.append([layout_name] + slide_ranges)
    
    # Convert the data into a pandas DataFrame
    max_cols = max(len(row) for row in layout_data)
    column_names = ["Layout Name"] + [f"Slides {i}" for i in range(1, max_cols)]
    df = pd.DataFrame(layout_data, columns=column_names).fillna("")
    
    # Save the DataFrame to an Excel file
    df.to_excel(output_excel_file, index=False, engine="openpyxl")
    print(f"Excel file saved to {output_excel_file}")
# Example usage
presentation_path = r"F:\5dmt Shashat\Codes and Files\Data\CopyData\قداس.pptx"  # Path to your PowerPoint presentation
output_excel_path = "layout_slides.xlsx"  # Path to save the Excel file
extract_slide_layouts_to_excel(presentation_path, output_excel_path)

# print(f"Layout and slide information has been saved to {output_excel_path}.")

# import win32com.client

# def link_textbox_to_custom_show(ppt_path, slide_number, textbox_name, custom_show_name):
#     ppMouseClick = 1
#     ppActionNamedSlideShow = 7

#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#     powerpoint.Visible = True

#     presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)

#     try:
#         slide = presentation.Slides(slide_number)

#         # Find the textbox by name
#         shape = None
#         for s in slide.Shapes:
#             if s.Name == textbox_name:
#                 shape = s
#                 break

#         if not shape:
#             print(f"[❌] Textbox named '{textbox_name}' not found on slide {slide_number}.")
#             return

#         action = shape.ActionSettings(ppMouseClick)

#         # Set action first
#         action.Action = ppActionNamedSlideShow
#         action.SlideShowName = custom_show_name
#         action.ShowAndReturn = True  # Only now is it safe to set

#         presentation.Save()

#     finally:
#         presentation.Close()
#         powerpoint.Quit()

# # Example usage
# link_textbox_to_custom_show(
#     ppt_path=r"F:\5dmt Shashat\Codes and Files\قداس.pptx",
#     slide_number=25,
#     textbox_name="TextBox 5",  # Exact shape name
#     custom_show_name="ابيكران"
# )

# from pptx import Presentation

# def replace_text_in_slides(pptx_path, old_text, new_text, start_slide, end_slide, output_path=None):
#     prs = Presentation(pptx_path)
    
#     # Ensure slide numbers are within bounds
#     total_slides = len(prs.slides)
#     start = max(1, start_slide)
#     end = min(end_slide, total_slides)
    
#     for i in range(start - 1, end):  # slide indices are 0-based
#         slide = prs.slides[i]
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         if old_text in run.text:
#                             run.text = run.text.replace(old_text, new_text)
    
#     # Save to the same file if no output_path is given
#     if output_path is None:
#         output_path = pptx_path
#     prs.save(output_path)
#     print(f"Text replacement completed in slides {start_slide} to {end_slide}.")

# # === Example usage ===
# replace_text_in_slides(
#     pptx_path=r"Data\القطمارس\الايام\القطمارس السنوي ايام.pptx",
#     old_text="الإنجيل من",
#     new_text="إنجيل باكر",
#     start_slide=0,
#     end_slide=1156,
#     output_path=r"Data\القطمارس\الايام\القطمارس السنوي ايام.pptx"
# )


# import win32com.client as win32

# def arrange_sections_alphabetically(pptx_path, start_section, end_section):
#     # Start PowerPoint
#     powerpoint = win32.Dispatch("PowerPoint.Application")
#     powerpoint.Visible = True
    
#     # Open the presentation
#     presentation = powerpoint.Presentations.Open(pptx_path)
    
#     sections = presentation.SectionProperties
#     section_count = sections.Count
    
#     print(f"Total sections: {section_count}")
    
#     # Build list of sections
#     section_list = []
#     for i in range(1, section_count + 1):  # COM is 1-based
#         section_name = sections.Name(i)
#         section_list.append({
#             "index": i,
#             "name": section_name,
#             "first_slide": sections.FirstSlide(i)
#         })
#         print(f"Section {i}: {section_name}")
    
#     # Find start and end indices
#     try:
#         start_idx = next(i for i, s in enumerate(section_list) if s["name"] == start_section)
#         end_idx = next(i for i, s in enumerate(section_list) if s["name"] == end_section)
#     except StopIteration:
#         print("❌ Start or End section not found!")
#         presentation.Close()
#         return
    
#     if start_idx > end_idx:
#         print("❌ Start section comes after End section!")
#         presentation.Close()
#         return
    
#     # Extract range
#     target_sections = section_list[start_idx:end_idx+1]
#     print(f"Sections to sort: {len(target_sections)}")
    
#     # Print original order
#     print("Original section order:")
#     for i, sec in enumerate(target_sections):
#         print(f"{i+1}. {sec['name']}")
    
#     # Sort sections by name using a basic string sort (may not be ideal for Arabic but let's try)
#     # We'll create a new list so we can track actual positions
#     sorted_sections = sorted(target_sections, key=lambda s: s["name"])
    
#     # Print the sorted order for debugging
#     print("Sorted section order:")
#     for i, sec in enumerate(sorted_sections):
#         print(f"{i+1}. {sec['name']}")
    
#     # We need to track section positions as they change during reordering
#     # So we'll create a function to get the current index of a section
#     def get_current_section_index(section_name):
#         for i in range(1, sections.Count + 1):
#             if sections.Name(i) == section_name:
#                 return i
#         return -1
    
#     # Reorder sections in PowerPoint
#     # Start from the beginning and place each section in order
#     target_pos = start_idx + 1  # 1-based position where we'll start placing sorted sections
    
#     for sec in sorted_sections:
#         # Get the current position of this section (which may have changed)
#         current_pos = get_current_section_index(sec["name"])
        
#         if current_pos != target_pos:
#             print(f"Moving section '{sec['name']}' from position {current_pos} to {target_pos}")
#             # Move the section
#             sections.Move(current_pos, target_pos)
        
#         target_pos += 1

# # Example usage:
# arrange_sections_alphabetically(
#     pptx_path=r"F:\5dmt Shashat\Codes and Files\Data\CopyData\كتاب المدائح.pptx",
#     start_section="رحلة الصوم الكبير",
#     end_section="كنيستي القبطية نشرت المسيحية"
# )

# import win32com.client

# def get_textboxes(slide):
#     """Return all textboxes (or placeholders) that contain text."""
#     boxes = []
#     for shape in slide.Shapes:
#         if shape.HasTextFrame:
#             if shape.TextFrame.HasText:
#                 if shape.TextFrame.TextRange.Text.strip() != "":
#                     boxes.append(shape)
#     return boxes

# def get_bottom_textbox(slide):
#     """Return the bottom-most textbox inside slide borders (largest Top value)."""
#     boxes = get_textboxes(slide)
#     if not boxes:
#         return None

#     slide_height = slide.Parent.PageSetup.SlideHeight  # total slide height
#     margin = 10  # tolerance in points

#     # Filter out shapes outside slide borders
#     inside_boxes = []
#     for s in boxes:
#         if (s.Top >= -margin and 
#             s.Left >= -margin and 
#             (s.Top + s.Height) <= (slide_height + margin)):
#             inside_boxes.append(s)

#     if not inside_boxes:
#         return None

#     return max(inside_boxes, key=lambda s: s.Top)

# def combine_textboxes(ppt_path):
#     app = win32com.client.Dispatch("PowerPoint.Application")
#     app.Visible = True
#     pres = app.Presentations.Open(ppt_path)

#     section_start = None
#     buffer_texts = []

#     for i, slide in enumerate(pres.Slides, start=1):
#         boxes = get_textboxes(slide)
#         print(f"Slide {i}: found {len(boxes)} textboxes")

#         # Blank slide → flush and reset
#         if len(boxes) == 0:
#             print("   -> Blank (reset section)")
#             if section_start and buffer_texts:
#                 main_box = get_bottom_textbox(section_start)
#                 if main_box:
#                     main_box.TextFrame.TextRange.Text += "\r\n" + "\r\n".join(buffer_texts)
#             buffer_texts = []
#             section_start = None
#             continue

#         # Slide with 2+ textboxes → new anchor
#         if len(boxes) >= 2:
#             print("   -> Multiple textboxes (new section anchor)")
#             # flush buffered text
#             if section_start and buffer_texts:
#                 main_box = get_bottom_textbox(section_start)
#                 if main_box:
#                     main_box.TextFrame.TextRange.Text += "\r\n" + "\r\n".join(buffer_texts)
#             buffer_texts = []
#             section_start = slide
#             continue

#         # Slide with exactly 1 textbox
#         if len(boxes) == 1:
#             if section_start is None:
#                 section_start = slide
#                 print("   -> Section start (1 textbox)")
#             else:
#                 print(f"   -> Cutting text from Slide {i} into Slide {section_start.SlideIndex}")
#                 buffer_texts.append(boxes[0].TextFrame.TextRange.Text)
#                 # CUT (clear text after saving)
#                 boxes[0].TextFrame.TextRange.Text = ""

#     # Final flush
#     if section_start and buffer_texts:
#         main_box = get_bottom_textbox(section_start)
#         if main_box:
#             main_box.TextFrame.TextRange.Text += "\r\n" + "\r\n".join(buffer_texts)

#     print("✅ Finished merging textboxes (cut + pasted into bottom textbox).")

# def clean_bottom_textbox(pptx_path):
#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#     powerpoint.Visible = True
    
#     presentation = powerpoint.Presentations.Open(pptx_path)

#     for slide in presentation.Slides:
#         bottom_shape = None
#         max_top = -1

#         shapes = slide.Shapes
#         for i in range(1, shapes.Count + 1):
#             shape = shapes.Item(i)
#             if shape.HasTextFrame and shape.TextFrame.HasText:
#                 if shape.Top > max_top:
#                     bottom_shape = shape
#                     max_top = shape.Top

#         if bottom_shape is not None:
#             txt = bottom_shape.TextFrame.TextRange.Text

#             # 1. Replace newlines with spaces
#             txt = txt.replace("\r\n", " ").replace("\r", " ").replace("\n", " ")

#             # 2. Remove multiple spaces
#             while "  " in txt:
#                 txt = txt.replace("  ", " ")

#             # 3. Remove leading and trailing spaces/newlines
#             txt = txt.strip()

#             bottom_shape.TextFrame.TextRange.Text = txt

# def delete_white_textboxes(pptx_path):
#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#     powerpoint.Visible = True
    
#     presentation = powerpoint.Presentations.Open(pptx_path)

#     for slide in presentation.Slides:
#         shapes = slide.Shapes
#         # Loop backwards because deleting while iterating forward causes skips
#         for i in range(shapes.Count, 0, -1):
#             shape = shapes.Item(i)
#             if shape.HasTextFrame and shape.TextFrame.HasText:
#                 font = shape.TextFrame.TextRange.Font
#                 if font.Color.RGB == 16777215:  # RGB(255,255,255) in decimal
#                     shape.Delete()

# def style_combined_textbox(textbox):
#     try:
#         text_range = textbox.TextFrame.TextRange
#         para_format = text_range.ParagraphFormat

#         # Remove underline from all text
#         text_range.Font.Underline = False

#         # Align text to right
#         para_format.Alignment = 3  # ppAlignRight

#         # Vertical center alignment (fixed)
#         textbox.TextFrame.VerticalAnchor = 3  # ppVerticalAnchorMiddle

#         # Turn on word wrap
#         textbox.TextFrame.WordWrap = True

#         # Ensure it's not autosizing (keep manual layout)
#         textbox.TextFrame.AutoSize = 0  # ppAutoSizeNone

#         # Make bullet list with star (★)
#         for paragraph in text_range.Paragraphs():
#             p_format = paragraph.ParagraphFormat
#             p_format.Bullet.Visible = True
#             p_format.Bullet.Font.Name = "Wingdings"
#             p_format.Bullet.Character = 118  # Unicode star
#             p_format.Bullet.RelativeSize = 1

#         textbox.Left = 0 * 28.35                       # 0 cm from left
#         textbox.Top = 12.07 * 28.35                    # 12.07 cm from top

#         textbox.Width = 25.4 * 28.35                   # 25.4 cm width
#         textbox.Height = 6.98 * 28.35                  # 6.98 cm height

#         print("✔ Styled combined textbox successfully.")

#     except Exception as e:
#         print(f"⚠️ Error styling textbox: {e}")

# def combine_textboxes_v2(ppt_path):
#     app = win32com.client.Dispatch("PowerPoint.Application")
#     app.Visible = True
#     pres = app.Presentations.Open(ppt_path)

#     section_start = None
#     buffer_texts = []

#     slides_to_delete = []

#     for i, slide in enumerate(list(pres.Slides), start=1):
#         boxes = get_textboxes(slide)
#         print(f"Slide {i}: found {len(boxes)} textboxes")

#         # Blank slide → flush and reset
#         if len(boxes) == 0:
#             print("   -> Blank (reset section)")
#             if section_start and buffer_texts:
#                 main_box = get_bottom_textbox(section_start)
#                 if main_box:
#                     main_box.TextFrame.TextRange.Text += "\r\n" + "\r\n".join(buffer_texts)
#                     style_combined_textbox(main_box)
#             buffer_texts = []
#             section_start = None
#             continue

#         # Slide with 2+ textboxes → new anchor
#         if len(boxes) >= 2:
#             print("   -> Multiple textboxes (new section anchor)")
#             # flush buffered text
#             if section_start and buffer_texts:
#                 main_box = get_bottom_textbox(section_start)
#                 if main_box:
#                     main_box.TextFrame.TextRange.Text += "\r\n" + "\r\n".join(buffer_texts)
#                     style_combined_textbox(main_box)
#             buffer_texts = []
#             section_start = slide
#             continue

#         # Slide with exactly 1 textbox
#         if len(boxes) == 1:
#             if section_start is None:
#                 section_start = slide
#                 print("   -> Section start (1 textbox)")
#             else:
#                 print(f"   -> Cutting text from Slide {i} into Slide {section_start.SlideIndex}")
#                 buffer_texts.append(boxes[0].TextFrame.TextRange.Text.strip())
#                 boxes[0].TextFrame.TextRange.Text = ""  # CUT text
#                 slides_to_delete.append(slide)  # mark for deletion later

#     # Final flush
#     if section_start and buffer_texts:
#         main_box = get_bottom_textbox(section_start)
#         if main_box:
#             main_box.TextFrame.TextRange.Text += "\r\n" + "\r\n".join(buffer_texts)
#             style_combined_textbox(main_box)

#     # Delete slides marked for deletion + hidden blank slides
#     for slide in reversed(list(pres.Slides)):
#         if slide in slides_to_delete:
#             print(f"Deleting slide {slide.SlideIndex} (merged text moved)")
#             slide.Delete()
#         elif slide.SlideShowTransition.Hidden:
#             boxes = get_textboxes(slide)
#             if len(boxes) == 0 or all(not b.TextFrame.TextRange.Text.strip() for b in boxes):
#                 print(f"Deleting hidden blank slide {slide.SlideIndex}")
#                 slide.Delete()

#     print("✅ Finished merging textboxes, formatting, and deleting empty slides.")

# def turn_bottom_text_white(pptx_path):
#     """Turn the text in the bottom textbox of each slide to white."""
#     powerpoint = win32com.client.Dispatch("PowerPoint.Application")
#     powerpoint.Visible = True

#     presentation = powerpoint.Presentations.Open(pptx_path)

#     for slide in presentation.Slides:
#         bottom_box = get_bottom_textbox(slide)
#         if not bottom_box:
#             continue

#         try:
#             text_range = bottom_box.TextFrame.TextRange

#             # 1. Set all text to white
#             text_range.Font.Color.RGB = 16777215  # RGB(255,255,255)

#             # 2. Set bullet color to yellow (#FFCC00)
#             bullet_rgb = int("00C0FF", 16)  # reversed byte order for COM (BGR)
#             # Explanation: PowerPoint COM uses BGR, not RGB, so #FFCC00 → 0x00CCFF
#             for paragraph in text_range.Paragraphs():
#                 if paragraph.ParagraphFormat.Bullet.Visible:
#                     paragraph.ParagraphFormat.Bullet.Font.Color.RGB = bullet_rgb

#             print(f"✔ Updated Slide {slide.SlideIndex}: text white, bullets yellow.")
#         except Exception as e:
#             print(f"⚠️ Error on Slide {slide.SlideIndex}: {e}")

# if __name__ == "__main__":
    # Change path to your presentation
    # ppt_file = r"F:\5dmt Shashat\Codes and Files\Data\القطمارس\السنكسار.pptx"
#     # clean_bottom_textbox(ppt_file)
#     combine_textboxes_v2(ppt_file)
#     # delete_white_textboxes(ppt_file)
    # turn_bottom_text_white(ppt_file)