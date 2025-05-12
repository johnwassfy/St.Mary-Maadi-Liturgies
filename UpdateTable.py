from datetime import timedelta
import os
from commonFunctions import relative_path, read_excel_cell, write_to_excel_cell, find_values_in_row
import asyncio
from copticDate import CopticCalendar
import pptx
from spire.presentation import Presentation
from openpyxl import load_workbook

file_path = relative_path(r"Tables.xlsx")
sheet_name = "المناسبات" 
copticDate = CopticCalendar()
copticYear = copticDate.gregorian_to_coptic()[0]

def a3yad ():

    # عيد الميلاد و برامون الميلاد و الغطاس
    if (copticDate.is_leap_year(copticYear-1)):
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F5", 28))
    else:
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F5", 29))
    YoumEidElmilad = read_excel_cell(file_path, sheet_name, "F5")
    if(copticDate.coptic_to_gregorian([copticYear,4, YoumEidElmilad]).weekday() == 7):
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F4", YoumEidElmilad-3))
    elif(copticDate.coptic_to_gregorian([copticYear,4, YoumEidElmilad]).weekday() == 6):
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F4", YoumEidElmilad-2))
    else:
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F4", YoumEidElmilad-1))

    if(copticDate.coptic_to_gregorian([copticYear, 5, 11]).weekday() == 7):
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F7", 8))
    elif(copticDate.coptic_to_gregorian([copticYear, 5, 11]).weekday() == 6):
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F7", 9))
    else:
        asyncio.run(write_to_excel_cell(file_path, sheet_name, "F7", 10))


    # ##حساب دور القمر
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "I25", (copticYear%19)-1))
    moonDay , moonMonth = find_values_in_row(file_path, sheet_name, "I", (copticYear%19)-1)
    gd = copticDate.coptic_to_gregorian([copticYear, moonMonth, moonDay])
    weekday = gd.weekday()
    days_until_sunday = (6 - weekday) % 7
    if days_until_sunday == 0:
            days_until_sunday = 7  # If today is Sunday, set days_until_sunday to 7

    nextSunday = gd + timedelta(days=days_until_sunday)

    rDate = copticDate.gregorian_to_coptic(nextSunday)
    rDay  = rDate[2]
    rMonth =  rDate[1]

    # # عيد القيامة
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F22", rDay))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D22", rMonth))

    # # احد الشعانين
    date = copticDate.coptic_date_before(7, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F18", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D18", date[1]))

    # # جمعة ختام الصوم
    date = copticDate.coptic_date_before(9, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F16", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D16", date[1]))

    # # سبت لعازر
    date = copticDate.coptic_date_before(8, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F17", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D17", date[1]))

    # # خميس العهد
    date = copticDate.coptic_date_before(3, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F19", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D19", date[1]))

    # # الجمعة العظيمة
    date = copticDate.coptic_date_before(2, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F20", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D20", date[1]))

    # # سبت النور
    date = copticDate.coptic_date_before(1, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F21", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D21", date[1]))

    # # عيد الصعود
    date = copticDate.coptic_date_after(39, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F24", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D24", date[1]))

    # # عيد العنصرة
    date = copticDate.coptic_date_after(49, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F25", date[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D25", date[1]))

    ## بداية الصوم الكبير
    FastingStartDate = copticDate.coptic_date_before(55, [copticYear, rMonth, rDay])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F13", FastingStartDate[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D13", FastingStartDate[1]))

    # ## حساب بداية صوم نينوى و فصح يونان
    NynowaStartDate = copticDate.coptic_date_before(14, [copticYear, FastingStartDate[1], FastingStartDate[2]])
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F11", NynowaStartDate[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D11", NynowaStartDate[1]))
    Fes7Younan = copticDate.coptic_date_after(3, NynowaStartDate)
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "F12", Fes7Younan[2]))
    asyncio.run(write_to_excel_cell(file_path, sheet_name, "D12", Fes7Younan[1]))


wb = load_workbook(file_path)
search_words = ["معلمنا بولس الرسول", "الكاثوليكون", "الإبركسيس", "المزمور (", "الإنجيل من"]
search_words2 = ["المزمور (", "الإنجيل من"]

def katamarsEl5amasyn():
    ws = wb["قطمارس الخماسين"]

    # Set column widths
    ws.column_dimensions['A'].width = 10
    ws.column_dimensions['B'].width = 15

    # Set ending day and month
    end_day = read_excel_cell(file_path, sheet_name, "F25")
    end_month = read_excel_cell(file_path, sheet_name, "D25")

    # Loop through 50 days in reverse
    for i in range(51, 1, -1):
        # Write day and month to the worksheet
        ws.cell(row=i, column=1).value = end_day
        ws.cell(row=i, column=2).value = end_month

        # Decrement day and adjust month if necessary
        end_day -= 1
        if end_day < 1:
            end_day = 30
            end_month -= 1

    pptx_file = relative_path(r"Data\القطمارس\قطمارس الخماسين (القداس).pptx")
    # Extract data from PowerPoint file
    prs = pptx.Presentation(pptx_file)
    matching_slides = {word: [] for word in search_words}
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        
        for word in search_words:
            if word.lower() in text:
                matching_slides[word].append(i)

    max_length = max(len(matching_slides[word]) for word in search_words)

    # Insert data into columns C to H for rows 2 to max_length + 1
    for i in range(max_length):
        slide_numbers = []
        for word in search_words:
            if i < len(matching_slides[word]):
                slide_numbers.append(matching_slides[word][i])
            else:
                slide_numbers.append('')
        
        for j, data in enumerate(slide_numbers, start=3):  # Start from column C
            ws.cell(row=i+2, column=j, value=data)

    # Copy values from column C (row 3 to row 50) to column H (row 2 to row 49)
    for row in range(3, 52):
        ws.cell(row=row - 1, column=8, value=ws.cell(row=row, column=3).value - 2 )

    # Insert total number of slides in the first cell of row 48 (column H)
    ws.cell(row=51, column=8, value=total_slides - 1)

    wb.save(file_path)

def katamarsOdasElsanawyAyam():
    ws = wb["القطمارس السنوي القداس"]

    for row in ws.iter_rows(min_row=2, min_col=1, max_col=ws.max_column, values_only=False):
        for cell in row:
            cell.value = None

    data = [
        (1, 1), (2, 1), (8, 1), (16, 1), (17, 1), (18, 1), (19, 1), (21, 1), (26, 1), (12, 2),
        (14, 2), (22, 2), (27, 2), (8, 3), (9, 3), (12, 3), (15, 3), (17, 3), (22, 3), (24, 3),
        (25, 3), (27, 3), (28, 3), (29, 3), (22, 4), (23,4), (28, 4), (29, 4), (30, 4), (1, 5), (3, 5),
        (4, 5), (6, 5), (10, 5), (11, 5), (12, 5), (13, 5), (22, 5), (26, 5), (30, 5), (2, 6),
        (13, 7), (29, 7), (23, 8), (27, 8), (30, 8), (1, 9), (10, 9), (20, 9), (24, 9), (26, 9),
        (2, 10), (16, 10), (30, 10), (3, 11), (5, 11), (20, 11), (3, 12), (13, 12), (17, 12),
        (25, 12), (26, 12), (28, 12), (29, 12), (30, 12), (1, 13), (2, 13), (3, 13), (4, 13), (6, 13)
    ]


    # Insert the data into the worksheet starting from row 2
    for i, row in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=row[0])
        ws.cell(row=i, column=2, value=row[1])

    pptx_file = relative_path(r"Data\القطمارس\الايام\القطمارس السنوي ايام (القداس).pptx")
    # Extract data from PowerPoint file
    prs = pptx.Presentation(pptx_file)
    prs2 = Presentation()
    prs2.LoadFromFile(pptx_file)
    matching_slides = {word: [] for word in search_words}
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        
        for word in search_words:
            if word.lower() in text:
                matching_slides[word].append(i)

    max_length = max(len(matching_slides[word]) for word in search_words)

    # Insert data into columns C to H for rows 2 to max_length + 1
    for i in range(max_length):
        slide_numbers = []
        for word in search_words:
            if i < len(matching_slides[word]):
                slide_numbers.append(matching_slides[word][i])
            else:
                slide_numbers.append('')
        
        for j, data in enumerate(slide_numbers, start=3):  # Start from column C
            ws.cell(row=i+2, column=j, value=data)
        
        for row in ws.iter_rows(min_row=2, min_col=7, max_col=8):
                g_cell = row[0]
                h_cell = row[1]

                # Read the value in column G
                start_slide = g_cell.value
                if start_slide is not None:
                    # Find the first not hidden slide after the start_slide
                    for slide_number in range(start_slide, prs2.Slides.Count):
                        if not prs2.Slides[slide_number].Hidden:
                            h_cell.value = slide_number  # Slides are 1-indexed in PowerPoint
                            break

    wb.save(file_path)

def katamarsOdasElsanawyA7ad():
    ws = wb["قطمارس الاحاد للقداس"]
    for row in ws.iter_rows(min_row=2, min_col=1, max_col=ws.max_column, values_only=False):
        for cell in row:
            cell.value = None
    data = [
        (1, 1), (2, 1), (3, 1), (4, 1), (1, 2), (2, 2), (3, 2), (4, 2), (1, 3), (2, 3), (3, 3), (4, 3),
        (1, 4), (2, 4), (3, 4), (4, 4), (1, 5), (2, 5), (3, 5), (4, 5), (1, 6), (2, 6), (3, 6), (4, 6),
        (3, 9), (4, 9), (1, 10), (2, 10), (3, 10), (4, 10), (1, 11), (2, 11), (3, 11), (4, 11), (1, 12), 
        (2, 12), (3, 12), (4, 12), (1, 13), (5, 1), (5, 2), (5, 3), (5, 4), (5, 5), (5, 6), (5, 9), 
        (5, 10), (5, 11), (5, 12)
    ]
    for i, row in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=row[0])
        ws.cell(row=i, column=2, value=row[1])

    pptx_file = relative_path(r"Data\القطمارس\الاحاد\القطمارس السنوي احاد (القداس).pptx")
    # Extract data from PowerPoint file
    prs = pptx.Presentation(pptx_file)
    matching_slides = {word: [] for word in search_words}
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        
        for word in search_words:
            if word.lower() in text:
                matching_slides[word].append(i)

    max_length = max(len(matching_slides[word]) for word in search_words)

    # Insert data into columns C to H for rows 2 to max_length + 1
    for i in range(max_length):
        slide_numbers = []
        for word in search_words:
            if i < len(matching_slides[word]):
                slide_numbers.append(matching_slides[word][i])
            else:
                slide_numbers.append('')
        
        for j, data in enumerate(slide_numbers, start=3):  # Start from column C
            ws.cell(row=i+2, column=j, value=data)

        for row in range(3, 42):
            original_value = ws.cell(row=row, column=3).value
            if original_value is not None:
                adjusted_value = original_value - 2
                ws.cell(row=row - 1, column=8, value=adjusted_value)

        # Insert total number of slides in the last cell of row 41 (column E)
        ws.cell(row=41, column=8, value=total_slides - 1)
        # Get the values from row 41, columns C to H
        values_to_copy = [ws.cell(row=41, column=col).value for col in range(3, 9)]
        # Paste the values into the next 9 rows
        for row in range(42, 51):
            for col_index, value in enumerate(values_to_copy, start=3):
                ws.cell(row=row, column=col_index, value=value)

    wb.save(file_path)

def katamars3ashyaElsanawyA7ad():
    ws = wb["قطمارس الاحاد للعشية"]
    for row in ws.iter_rows(min_row=2, min_col=1, max_col=ws.max_column, values_only=False):
        for cell in row:
            cell.value = None
    data = [
        (1, 1), (2, 1), (3, 1), (4, 1), (1, 2), (2, 2), (3, 2), (4, 2), (1, 3), (2, 3), (3, 3), (4, 3),
        (1, 4), (2, 4), (3, 4), (4, 4), (1, 5), (2, 5), (3, 5), (4, 5), (1, 6), (2, 6), (3, 6), (4, 6),
        (3, 9), (4, 9), (1, 10), (2, 10), (3, 10), (4, 10), (1, 11), (2, 11), (3, 11), (4, 11), (1, 12), 
        (2, 12), (3, 12), (4, 12), (1, 13), (5, 1), (5, 2), (5, 3), (5, 4), (5, 5), (5, 6), (5, 9), 
        (5, 10), (5, 11), (5, 12)
    ]
    for i, row in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=row[0])
        ws.cell(row=i, column=2, value=row[1])

    pptx_file = relative_path(r"Data\القطمارس\الاحاد\القطمارس السنوي احاد (عشية).pptx")
    # Extract data from PowerPoint file
    prs = pptx.Presentation(pptx_file)
    matching_slides = {word: [] for word in search_words2}
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        
        for word in search_words2:
            if word.lower() in text:
                matching_slides[word].append(i)

    max_length = max(len(matching_slides[word]) for word in search_words2)
    # Insert data into columns C to H for rows 2 to max_length + 1
    for i in range(max_length):
        slide_numbers = []
        for word in search_words2:
            if i < len(matching_slides[word]):
                slide_numbers.append(matching_slides[word][i])
            else:
                slide_numbers.append('')
        
        for j, data in enumerate(slide_numbers, start=3):  # Start from column C
            ws.cell(row=i+2, column=j, value=data)

        for row in range(3, 42):
            original_value = ws.cell(row=row, column=3).value
            if original_value is not None:
                adjusted_value = original_value - 2
                ws.cell(row=row - 1, column=5, value=adjusted_value)

        # Insert total number of slides in the last cell of row 41 (column E)
        ws.cell(row=41, column=5, value=total_slides - 1)
        # Get the values from row 41, columns C to H
        values_to_copy = [ws.cell(row=41, column=col).value for col in range(3, 9)]
        # Paste the values into the next 9 rows
        for row in range(42, 51):
            for col_index, value in enumerate(values_to_copy, start=3):
                ws.cell(row=row, column=col_index, value=value)

        wb.save(file_path)

def katamarsBakerElsanawyA7ad():
    ws = wb["قطمارس الاحاد لباكر"]
    for row in ws.iter_rows(min_row=2, min_col=1, max_col=ws.max_column, values_only=False):
        for cell in row:
            cell.value = None
    data = [
        (1, 1), (2, 1), (3, 1), (4, 1), (1, 2), (2, 2), (3, 2), (4, 2), (1, 3), (2, 3), (3, 3), (4, 3),
        (1, 4), (2, 4), (3, 4), (4, 4), (1, 5), (2, 5), (3, 5), (4, 5), (1, 6), (2, 6), (3, 6), (4, 6),
        (3, 9), (4, 9), (1, 10), (2, 10), (3, 10), (4, 10), (1, 11), (2, 11), (3, 11), (4, 11), (1, 12), 
        (2, 12), (3, 12), (4, 12), (1, 13), (5, 1), (5, 2), (5, 3), (5, 4), (5, 5), (5, 6), (5, 9), 
        (5, 10), (5, 11), (5, 12)
    ]
    for i, row in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=row[0])
        ws.cell(row=i, column=2, value=row[1])

    pptx_file = relative_path(r"Data\القطمارس\الاحاد\القطمارس السنوي احاد (باكر).pptx")
    # Extract data from PowerPoint file
    prs = pptx.Presentation(pptx_file)
    matching_slides = {word: [] for word in search_words2}
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        
        for word in search_words2:
            if word.lower() in text:
                matching_slides[word].append(i)

    max_length = max(len(matching_slides[word]) for word in search_words2)
    # Insert data into columns C to H for rows 2 to max_length + 1
    for i in range(max_length):
        slide_numbers = []
        for word in search_words2:
            if i < len(matching_slides[word]):
                slide_numbers.append(matching_slides[word][i])
            else:
                slide_numbers.append('')
        
        for j, data in enumerate(slide_numbers, start=3):  # Start from column C
            ws.cell(row=i+2, column=j, value=data)

        for row in range(3, 42):
            original_value = ws.cell(row=row, column=3).value
            if original_value is not None:
                adjusted_value = original_value - 2
                ws.cell(row=row - 1, column=5, value=adjusted_value)

        # Insert total number of slides in the last cell of row 41 (column E)
        ws.cell(row=41, column=5, value=total_slides - 1)
        # Get the values from row 41, columns C to H
        values_to_copy = [ws.cell(row=41, column=col).value for col in range(3, 9)]
        # Paste the values into the next 9 rows
        for row in range(42, 51):
            for col_index, value in enumerate(values_to_copy, start=3):
                ws.cell(row=row, column=col_index, value=value)

        wb.save(file_path)

def katamars3ashyaElsanawyAyam():
    ws = wb["القطمارس السنوي العشية"]

    for row in ws.iter_rows(min_row=2, min_col=1, max_col=ws.max_column, values_only=False):
        for cell in row:
            cell.value = None

    data = [
        (1, 1), (2, 1), (8, 1), (16, 1), (17, 1), (18, 1), (19, 1), (21, 1), (26, 1), (12, 2),
        (14, 2), (22, 2), (27, 2), (8, 3), (9, 3), (12, 3), (15, 3), (17, 3), (22, 3), (24, 3),
        (25, 3), (27, 3), (28, 3), (29, 3), (22, 4), (23,4), (28, 4), (29, 4), (30, 4), (1, 5), (3, 5),
        (4, 5), (6, 5), (10, 5), (11, 5), (12, 5), (13, 5), (22, 5), (26, 5), (30, 5), (2, 6),
        (13, 7), (29, 7), (23, 8), (27, 8), (30, 8), (1, 9), (10, 9), (20, 9), (24, 9), (26, 9),
        (2, 10), (16, 10), (30, 10), (3, 11), (5, 11), (20, 11), (3, 12), (13, 12), (17, 12),
        (25, 12), (26, 12), (28, 12), (29, 12), (30, 12), (1, 13), (2, 13), (3, 13), (4, 13), (6, 13)
    ]


    for i, row in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=row[0])
        ws.cell(row=i, column=2, value=row[1])

    pptx_file = relative_path(r"Data\القطمارس\الايام\القطمارس السنوي ايام (عشية).pptx")
    # Extract data from PowerPoint file
    prs = pptx.Presentation(pptx_file)
    prs2 = Presentation()
    prs2.LoadFromFile(pptx_file)
    matching_slides = {word: [] for word in search_words2}
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        
        for word in search_words2:
            if word.lower() in text:
                matching_slides[word].append(i)

    max_length = max(len(matching_slides[word]) for word in search_words2)

    # Insert data into columns C to H for rows 2 to max_length + 1
    for i in range(max_length):
        slide_numbers = []
        for word in search_words2:
            if i < len(matching_slides[word]):
                slide_numbers.append(matching_slides[word][i])
            else:
                slide_numbers.append('')
        
        for j, data in enumerate(slide_numbers, start=3):  # Start from column C
            ws.cell(row=i+2, column=j, value=data)

        for row in ws.iter_rows(min_row=2, min_col=4, max_col=5):
                g_cell = row[0]
                h_cell = row[1]

                # Read the value in column G
                start_slide = g_cell.value
                if start_slide is not None:
                    # Find the first not hidden slide after the start_slide
                    for slide_number in range(start_slide, prs2.Slides.Count):
                        if not prs2.Slides[slide_number].Hidden:
                            h_cell.value = slide_number  # Slides are 1-indexed in PowerPoint
                            break

    wb.save(file_path)

def katamarsBakerElsanawyAyam():
    ws = wb["القطمارس السنوي باكر"]

    for row in ws.iter_rows(min_row=2, min_col=1, max_col=ws.max_column, values_only=False):
        for cell in row:
            cell.value = None

    data = [
        (1, 1), (2, 1), (8, 1), (16, 1), (17, 1), (18, 1), (19, 1), (21, 1), (26, 1), (12, 2),
        (14, 2), (22, 2), (27, 2), (8, 3), (9, 3), (12, 3), (15, 3), (17, 3), (22, 3), (24, 3),
        (25, 3), (27, 3), (28, 3), (29, 3), (22, 4), (23,4), (28, 4), (29, 4), (30, 4), (1, 5), (3, 5),
        (4, 5), (6, 5), (10, 5), (11, 5), (12, 5), (13, 5), (22, 5), (26, 5), (30, 5), (2, 6),
        (13, 7), (29, 7), (23, 8), (27, 8), (30, 8), (1, 9), (10, 9), (20, 9), (24, 9), (26, 9),
        (2, 10), (16, 10), (30, 10), (3, 11), (5, 11), (20, 11), (3, 12), (13, 12), (17, 12),
        (25, 12), (26, 12), (28, 12), (29, 12), (30, 12), (1, 13), (2, 13), (3, 13), (4, 13), (6, 13)
    ]


    for i, row in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=row[0])
        ws.cell(row=i, column=2, value=row[1])

    pptx_file = relative_path(r"Data\القطمارس\الايام\القطمارس السنوي ايام (باكر).pptx")
    # Extract data from PowerPoint file
    prs = pptx.Presentation(pptx_file)
    prs2 = Presentation()
    prs2.LoadFromFile(pptx_file)
    matching_slides = {word: [] for word in search_words2}
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        
        for word in search_words2:
            if word.lower() in text:
                matching_slides[word].append(i)

    max_length = max(len(matching_slides[word]) for word in search_words2)

    # Insert data into columns C to H for rows 2 to max_length + 1
    for i in range(max_length):
        slide_numbers = []
        for word in search_words2:
            if i < len(matching_slides[word]):
                slide_numbers.append(matching_slides[word][i])
            else:
                slide_numbers.append('')
        
        for j, data in enumerate(slide_numbers, start=3):  # Start from column C
            ws.cell(row=i+2, column=j, value=data)

        for row in ws.iter_rows(min_row=2, min_col=4, max_col=5):
                g_cell = row[0]
                h_cell = row[1]

                # Read the value in column G
                start_slide = g_cell.value
                if start_slide is not None:
                    # Find the first not hidden slide after the start_slide
                    for slide_number in range(start_slide, prs2.Slides.Count):
                        if not prs2.Slides[slide_number].Hidden:
                            h_cell.value = slide_number  # Slides are 1-indexed in PowerPoint
                            break

    wb.save(file_path)

def Younan():
    ws = wb["صوم نينوى و فصح يونان"]

    # 1) Set starting day and month in columns A and B
    startday = read_excel_cell(file_path, sheet_name, "F11")
    startmonth = read_excel_cell(file_path, sheet_name, "D11")

    for i in range(1, 5):
        ws.cell(row=i + 1, column=1).value = startday   # Column A
        ws.cell(row=i + 1, column=2).value = startmonth # Column B
        startday += 1
        if startday > 30:
            startday = 1
            startmonth += 1

    # 2) Read the PPTX file
    pptx_file = relative_path(r"Data\القطمارس\الصوم الكبير و صوم نينوى\قرائات صوم نينوى و فصح يونان.pptx")
    prs = pptx.Presentation(pptx_file)

    # 3) Define search words and collect matching slides
    younan_search_words = [
        "مزمور عشية",
        "إنجيل عشية",
        "نبوات",
        "مزمور باكر",
        "إنجيل باكر",
        "معلمنا بولس الرسول",
        "الكاثوليكون",
        "الإبركسيس",
        "المزمور (",
        "الإنجيل من"
    ]
    matching_slides = {word: [] for word in younan_search_words}

    total_slides = len(prs.slides)
    
    # 4) Get all visible slide numbers
    visible_slides = [i + 1 for i, slide in enumerate(prs.slides) if slide._element.get("show") != "0"]

    # 5) Collect slide numbers for each search word
    for i, slide in enumerate(prs.slides, start=1):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        for word in younan_search_words:
            if word.lower() in text:
                matching_slides[word].append(i)

    # 6) Determine the maximum number of rows we need
    max_length = max(len(matching_slides[word]) for word in younan_search_words)

    # 7) Map each of the 10 search words to its target column, skipping E(5) and I(9)
    #    So the columns in order: C(3), D(4), F(6), G(7), H(8), J(10), K(11), L(12), M(13), N(14)
    target_columns = [3, 4, 6, 7, 8, 10, 11, 12, 13, 14]

    # 8) Insert data for each row
    for row_idx in range(max_length):
        # Gather this row's 10 slides
        row_slide_numbers = []
        for word in younan_search_words:
            if row_idx < len(matching_slides[word]):
                row_slide_numbers.append(matching_slides[word][row_idx])
            else:
                row_slide_numbers.append('')

        # Insert them into the correct columns
        for col_idx, slide_num in enumerate(row_slide_numbers):
            ws.cell(row=row_idx + 2, column=target_columns[col_idx], value=slide_num)

    # 9) For columns E(5), I(9), and O(15):
    #    - Read the slide from the "previous" column
    #    - Find the next visible slide after that number
    #    - Subtract 1 and write to E, I, or O
    last_slide = total_slides

    # We'll loop from row 2 through row max_length+1 (since data insertion starts at row 2)
    for row in range(2, max_length + 2):
        # a) Column E = next visible slide after the value in column D (4), minus 1
        col_d_value = ws.cell(row=row, column=4).value
        if col_d_value is not None and isinstance(col_d_value, int):
            next_slide = next((s for s in visible_slides if s > col_d_value), last_slide) - 1
            ws.cell(row=row, column=5).value = next_slide  # Column E

        # b) Column I = next visible slide after the value in column H (8), minus 1
        col_h_value = ws.cell(row=row, column=8).value
        if col_h_value is not None and isinstance(col_h_value, int):
            next_slide = next((s for s in visible_slides if s > col_h_value), last_slide) - 1
            ws.cell(row=row, column=9).value = next_slide  # Column I

        # c) Column O = next visible slide after the value in column N (14), minus 1
        col_n_value = ws.cell(row=row, column=14).value
        if col_n_value is not None and isinstance(col_n_value, int):
            next_slide = next((s for s in visible_slides if s > col_n_value), last_slide) - 1
            ws.cell(row=row, column=15).value = next_slide  # Column O

    # 10) Finally, save the workbook
    wb.save(file_path)

def ElsomElkbyr():
    ws = wb["قطمارس الصوم الكبير"]

    # Set column widths for reference
    ws.column_dimensions['A'].width = 10
    ws.column_dimensions['B'].width = 15

    # 1) Set starting day and month in columns A and B
    ElsomElKbyrday = read_excel_cell(file_path, sheet_name, "F13")
    ElsomElKbyrmonth = read_excel_cell(file_path, sheet_name, "D13")

    SbtElrefa3 = copticDate.coptic_date_before(2, [copticYear, ElsomElKbyrmonth, ElsomElKbyrday])
    day = SbtElrefa3[2]
    month = SbtElrefa3[1]

    # Loop through 50 days (from سبت الرفاع to سبت لعازر)
    for i in range(1, 51):
        ws.cell(row=i + 1, column=1).value = day     # Column A
        ws.cell(row=i + 1, column=2).value = month     # Column B
        day += 1
        if day > 30:
            day = 1
            month += 1

    # 2) Read the PPTX file
    pptx_file = relative_path(r"Data\القطمارس\الصوم الكبير و صوم نينوى\قطمارس الصوم الكبير.pptx")
    prs = pptx.Presentation(pptx_file)

    # 3) Define search words and collect matching slides
    younan_search_words = [
        "مزمور عشية",
        "إنجيل عشية",
        "نبوات",
        "مزمور باكر",
        "إنجيل باكر",
        "معلمنا بولس الرسول",
        "الكاثوليكون",
        "الإبركسيس",
        "المزمور (",
        "الإنجيل من",
        "مزمور مساء الاحد",
        "انجيل مساء الاحد"
    ]
    matching_slides = {word: [] for word in younan_search_words}

    total_slides = len(prs.slides)
    
    # 4) Get all visible slide numbers (non-hidden)
    visible_slides = [i + 1 for i, slide in enumerate(prs.slides) if slide._element.get("show") != "0"]

    # 5) Populate matching_slides with slide numbers for each search word
    for i, slide in enumerate(prs.slides, start=1):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        for word in younan_search_words:
            if word.lower() in text:
                matching_slides[word].append(i)

    # -------------------------------------------------------------------------
    # 6) Place slide numbers into the sheet according to your patterns.
    # -------------------------------------------------------------------------

    # A) For "مزمور عشية" and "إنجيل عشية":
    #    They are entered as pairs (in the same row) in columns C and D.
    #    The row sequence for the pairs is: 2, 3, 10, 17, 24, 31, ...
    mosmor_list = matching_slides["مزمور عشية"]
    engeel_list = matching_slides["إنجيل عشية"]
    num_pairs = max(len(mosmor_list), len(engeel_list))
    for i in range(num_pairs):
        if i == 0:
            row = 2
        elif i == 1:
            row = 3
        else:
            row = 3 + (i - 1) * 7
        # Place "مزمور عشية" in Column C (3)
        if i < len(mosmor_list):
            ws.cell(row=row, column=3).value = mosmor_list[i]
        # Place "إنجيل عشية" in Column D (4)
        if i < len(engeel_list):
            ws.cell(row=row, column=4).value = engeel_list[i]

    # B) For "نبوات" -> Column F (6)
    #    Pattern: start at row 4, enter 5 consecutive rows, skip 2 rows repeatedly,
    #    and then, before placing the last 6 slides, skip 2 rows, and then enter them consecutively.
    nboaat_list = matching_slides["نبوات"]
    n_nboaat = len(nboaat_list)
    if n_nboaat <= 6:
        row = 4
        for slide_num in nboaat_list:
            ws.cell(row=row, column=6).value = slide_num
            row += 1
    else:
        row = 4
        placed = 0
        # Place slides in chunks of 5 with a 2-row gap until only 6 remain.
        while placed < (n_nboaat - 6):
            chunk_size = min(5, (n_nboaat - 6) - placed)
            for i in range(chunk_size):
                ws.cell(row=row, column=6).value = nboaat_list[placed + i]
                row += 1
            placed += chunk_size
            if placed < (n_nboaat - 6):
                row += 2  # Skip 2 rows before next chunk

        # Before placing the last 6 slides, skip 2 rows.
        row += 2
        for i in range(n_nboaat - 6, n_nboaat):
            ws.cell(row=row, column=6).value = nboaat_list[i]
            row += 1

    # C) For "مزمور مساء الاحد" and "انجيل مساء الاحد":
    #    They are entered in columns P (16) and Q (17) respectively.
    #    Pattern: starting at row 3, then each pair is placed 7 rows down.
    mosmor_masaa_list = matching_slides["مزمور مساء الاحد"]
    for i, slide_num in enumerate(mosmor_masaa_list):
        row = 3 + i * 7
        ws.cell(row=row, column=16).value = slide_num  # Column P

    engeel_masaa_list = matching_slides["انجيل مساء الاحد"]
    for i, slide_num in enumerate(engeel_masaa_list):
        row = 3 + i * 7
        ws.cell(row=row, column=17).value = slide_num  # Column Q

    # D) The rest of the words are entered normally in consecutive rows starting at row 2,
    #    in their respective columns. (Skipped columns: E, I, O, and R)
    normal_columns = {
        "مزمور باكر": 7,            # Column G
        "إنجيل باكر": 8,            # Column H
        "معلمنا بولس الرسول": 10,    # Column J
        "الكاثوليكون": 11,         # Column K
        "الإبركسيس": 12,           # Column L
        "المزمور (": 13,           # Column M
        "الإنجيل من": 14            # Column N
    }
    for word, col in normal_columns.items():
        slides_list = matching_slides[word]
        row = 2
        for slide_num in slides_list:
            ws.cell(row=row, column=col).value = slide_num
            row += 1

    # -------------------------------------------------------------------------
    # 7) Now fill in the skipped columns by using the previous column's value.
    #    For each skipped column, if the previous column (which contains a slide number)
    #    is not empty and is an integer, we find the next visible slide after that number,
    #    subtract 1, and write it in the skipped column.
    #
    # Skipped columns and their corresponding previous columns:
    #   Column E (5)  <- Column D (4)  ("إنجيل عشية")
    #   Column I (9)  <- Column H (8)  ("إنجيل باكر")
    #   Column O (15) <- Column N (14) ("الإنجيل من")
    #   Column R (18) <- Column Q (17) ("انجيل مساء الاحد")
    # -------------------------------------------------------------------------

    last_slide = total_slides

    # For Column E: based on Column D
    for row in range(2, ws.max_row + 1):
        prev_val = ws.cell(row=row, column=4).value
        if isinstance(prev_val, int):
            next_slide = next((s for s in visible_slides if s > prev_val), last_slide) - 1
            ws.cell(row=row, column=5).value = next_slide

    # For Column I: based on Column H
    for row in range(2, ws.max_row + 1):
        prev_val = ws.cell(row=row, column=8).value
        if isinstance(prev_val, int):
            next_slide = next((s for s in visible_slides if s > prev_val), last_slide) - 1
            ws.cell(row=row, column=9).value = next_slide

    # For Column O: based on Column N
    for row in range(2, ws.max_row + 1):
        prev_val = ws.cell(row=row, column=14).value
        if isinstance(prev_val, int):
            next_slide = next((s for s in visible_slides if s > prev_val), last_slide) - 1
            ws.cell(row=row, column=15).value = next_slide

    # For Column R: based on Column Q
    for row in range(2, ws.max_row + 1):
        prev_val = ws.cell(row=row, column=17).value
        if isinstance(prev_val, int):
            next_slide = next((s for s in visible_slides if s > prev_val), last_slide) - 1
            ws.cell(row=row, column=18).value = next_slide

    # 8) Finally, save the workbook
    wb.save(file_path)

def Elsh3anyn():
    ws = wb["قرائات أحد الشعانين"]
    
    # 1) Set starting day and month in columns A and B
    Elsh3anynDay = read_excel_cell(file_path, sheet_name, "F18")
    Elsh3anynMonth = read_excel_cell(file_path, sheet_name, "D18")

    ws.cell(2, column=1).value = Elsh3anynDay     # Column A
    ws.cell(2, column=2).value = Elsh3anynMonth     # Column B

    # 2) Read the PPTX file
    pptx_file = relative_path(r"Data\القطمارس\قرائات احد الشعانين.pptx")
    prs = pptx.Presentation(pptx_file)

    # 3) Define search words and collect matching slides
    sha3anyn_search_words = [
        "مزمور عشية",
        "انجيل عشية",
        "مزمور باكر",
        "انجيل باكر",
        "معلمنا بولس الرسول",
        "الكاثوليكون",
        "الإبركسيس",
        "المزمور الأول",
        "الإنجيل الأول",
        "الإنجيل الثانى",
        "الإنجيل الثالث",
        "المزمور الثاني",
        "الإنجيل الرابع",
    ]
    matching_slides = {word: [] for word in sha3anyn_search_words}

    total_slides = len(prs.slides)

    # 4) Collect slide numbers for each search word
    for i, slide in enumerate(prs.slides, start=1):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text.lower()
        for word in sha3anyn_search_words:
            if word.lower() in text:
                matching_slides[word].append(i)

    # 5) Determine the maximum number of rows we need
    max_length = max(len(matching_slides[word]) for word in sha3anyn_search_words)

    # 6) Map each of the 10 search words to its target column, skipping E(5), H(8),  N(14), P(16), R(18), U(21)
    #    So the columns in order: C(3), D(4), F(6), G(7), I(9), J(10), K(11), L(12), M(13), O(15), Q(17), S(19), T(20)
    target_columns = [3, 4, 6, 7, 9, 10, 11, 12, 13, 15, 17, 19, 20]

    # 8) Insert data for each row
    for row_idx in range(max_length):
        # Gather this row's 10 slides
        row_slide_numbers = []
        for word in sha3anyn_search_words:
            if row_idx < len(matching_slides[word]):
                row_slide_numbers.append(matching_slides[word][row_idx])
            else:
                row_slide_numbers.append('')

        # Insert them into the correct columns
        for col_idx, slide_num in enumerate(row_slide_numbers):
            ws.cell(row=row_idx + 2, column=target_columns[col_idx], value=slide_num)
    
    # 9) For columns E(5), H(8), N(14), P(16), R(18), U(21):
    #    - Read the slide from the "next" column (F, I, O, Q, S, V respectively)
    #    - Subtract 1 and write to E, H, N, P, R, or U
    
    # Define the column pairs: (target_column, source_column)
    column_pairs = [
        (5, 6),   # E gets value from F minus 1
        (8, 9),    # H gets value from I minus 1
        (14, 15),  # N gets value from O minus 1
        (16, 17),  # P gets value from Q minus 1
        (18, 19),  # R gets value from S minus 1
    ]
    
    # Iterate through each row that has data (from row 2 to max_length + 1)
    for row_idx in range(2, max_length + 2):
        for target_col, source_col in column_pairs:
            source_value = ws.cell(row=row_idx, column=source_col).value
            if source_value and str(source_value).isdigit():
                # Subtract 1 only if the source value is a number
                ws.cell(row=row_idx, column=target_col, value=int(source_value) - 1)
            else:
                # If source is empty or not a number, leave target empty
                ws.cell(row=row_idx, column=target_col, value='')
    
    ws.cell(row=2, column=21, value=total_slides)
    
    # 10) Finally, save the workbook
    wb.save(file_path)

def All(progress_callback=None):
    total_steps = 9  # Total number of functions called
    current_step = 0

    Younan()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    ElsomElkbyr()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    Elsh3anyn()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    katamarsEl5amasyn()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    katamarsOdasElsanawyAyam()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    katamarsOdasElsanawyA7ad()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    katamars3ashyaElsanawyA7ad()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    katamars3ashyaElsanawyAyam()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))

    katamarsBakerElsanawyAyam()
    current_step += 1
    if progress_callback:
        progress_callback(int((current_step / total_steps) * 100))