import os
import sys
from openpyxl import load_workbook, utils
from pptx import Presentation
import win32com.client

def relative_path(relative_path):
    """Return absolute path to resource, for dev and PyInstaller"""
    base_path = getattr(sys, '_MEIPASS', os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(base_path, relative_path)

def load_background_image(background_label, image_path=r"Data\الصور\background2.png"):
    from PyQt5.QtGui import QPixmap
    pixmap = QPixmap(relative_path(image_path))
    if background_label:  # Check if the background label exists
        background_label.setPixmap(pixmap)
        background_label.setScaledContents(True)
        background_label.setObjectName("background_label")  # Set object name for the background label

def find_slide_num(excel_path, sheet_name, word, col_num):
    try:
        # Load the Excel workbook
        workbook = load_workbook(excel_path, read_only=True)
        
        # Select the worksheet by name
        worksheet = workbook[sheet_name]
        
        # Iterate over rows in the worksheet
        for row in worksheet.iter_rows(values_only=True):
            if row[0] == word:  # Assuming the words are in the first column
                return row[col_num]  # Return the corresponding number from col_num
        
        # If the word is not found, return a message indicating that
        return f"No corresponding number found for '{word}'."
    
    except Exception as e:
        # Handle any errors that might occur (e.g., file not found, sheet not found, etc.)
        return f"Error: {str(e)}"

def find_slide_num_v2(excel_path, sheet_name, word, search_col, offset):
    try:
        # Load the Excel workbook
        workbook = load_workbook(excel_path, read_only=True)
        
        # Select the worksheet by name
        worksheet = workbook[sheet_name]
        
        # Adjust search_col to zero-based index
        search_col -= 1
        
        for row in worksheet.iter_rows(values_only=True):
            if row[search_col] == word: 
                # Calculate the target column index
                target_col = search_col + offset
                if 0 <= target_col < len(row):
                    return row[target_col]  # Return the value directly
                else:
                    return f"Offset {offset} out of bounds for row with word '{word}'."
        
        return f"No corresponding number found for '{word}' in column {search_col + 1}."
    
    except Exception as e:
        return f"Error: {str(e)}"

def find_slide_nums_arrays(excel_path, sheet_name, words, col_nums):
    try:
        # Load the Excel workbook
        workbook = load_workbook(excel_path, read_only=True)
        
        # Select the worksheet by name
        worksheet = workbook[sheet_name]
        
        # Initialize an empty list to hold the results
        results = []
        
        # Iterate over the words and their corresponding column numbers
        for word, col_num in zip(words, col_nums):
            found = False
            # Iterate over rows in the worksheet
            for row in worksheet.iter_rows(values_only=True):
                if row[0] == word: 
                    # Assuming the words are in the first column
                    results.append(row[col_num])  # Append the result for the current word
                    found = True
                    break
            if not found:
                results.append(f"No corresponding number found for '{word}'.")
        
        return results
    
    except Exception as e:
        # Handle any errors that might occur (e.g., file not found, sheet not found, etc.)
        return f"Error: {str(e)}"

def find_slide_nums_arrays_v2(excel_path, sheet_name, words, search_col, offsets):
    try:
        # Load the Excel workbook
        workbook = load_workbook(excel_path, read_only=True)
        
        # Select the worksheet by name
        worksheet = workbook[sheet_name]
        
        # Initialize an empty list to hold the results
        results = []
        
        # Adjust search_col to zero-based index
        search_col -= 1
        
        # Iterate over the words and their corresponding offsets
        for word, offset in zip(words, offsets):
            found = False
            # Iterate over rows in the worksheet
            for row in worksheet.iter_rows(values_only=True):
                if row[search_col] == word: 
                    # Calculate the target column index
                    target_col = search_col + offset
                    if 0 <= target_col < len(row):
                        results.append(row[target_col])  # Append the value from the target column
                    else:
                        results.append(f"Offset {offset} out of bounds for row with word '{word}'.")
                    found = True
                    break
            if not found:
                results.append(f"No corresponding number found for '{word}' in column {search_col + 1}.")
        
        return results
    
    except Exception as e:
        # Handle any errors that might occur (e.g., file not found, sheet not found, etc.)
        return f"Error: {str(e)}"

def read_excel_cell(file_path, sheet_name, cell_address):
    # Load the workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    
    # Select the worksheet
    ws = wb[sheet_name]
    
    # Read the value from the specified cell
    cell_value = ws[cell_address].value
    
    # Close the workbook
    wb.close()
    
    return cell_value

def find_values_in_row(file_path, sheet_name, column_letter, search_value):
    # Load the workbook and select the sheet
    workbook = load_workbook(file_path)
    sheet = workbook[sheet_name]
    
    # Convert the column letter to a number (1-indexed)
    column_index = utils.column_index_from_string(column_letter)
    
    # Iterate through the rows in the specified column
    for row in sheet.iter_rows(min_col=column_index, max_col=column_index):
        for cell in row:
            if cell.value == search_value:
                # Get the row number
                row_num = cell.row
                # Get the values in the next two columns (J and K if column is I)
                value1 = sheet.cell(row=row_num, column=column_index + 1).value
                value2 = sheet.cell(row=row_num, column=column_index + 2).value
                return value1, value2
    
    # If the value is not found, return None
    return None, None

async def write_to_excel_cell(file_path, sheet_name, cell_address, value):
    # Load the workbook
    wb = load_workbook(filename=file_path)
    
    # Select the worksheet
    ws = wb[sheet_name]
    
    # Write the value to the specified cell
    ws[cell_address] = value
    
    # Save the workbook
    wb.save(file_path)
    
    # Close the workbook
    wb.close()

def read_column(file_path, sheet_name, column):
    # Load the Excel workbook
    wb = load_workbook(filename=file_path)
    
    # Select the worksheet
    ws = wb[sheet_name]
    
    # Initialize a list to store values
    column_values = []
    
    # Iterate over cells in the column until an empty cell is encountered
    for cell in ws[column]:
        if cell.value is not None:
            column_values.append(cell.value)
        else:
            # Break the loop if an empty cell is encountered
            break
    
    return column_values

def find_and_save_values(file_path, sheet_name, column_letter, search_value):
    # Load the workbook
    wb = load_workbook(file_path)
    
    # Select the worksheet
    ws = wb[sheet_name]

    # Get the maximum row count in the column
    max_row = ws.max_row

    # Iterate over the specified column
    for row in range(1, max_row + 1):
        cell_value = ws[column_letter + str(row)].value
        if cell_value == search_value:
            # Save corresponding values from adjacent columns
            value1 = ws.cell(row=row, column=column_letter_to_number(column_letter) + 1).value
            value2 = ws.cell(row=row, column=column_letter_to_number(column_letter) + 2).value
            return value1, value2

    # If the search value is not found, return None
    return None, None

def column_letter_to_number(letter):
    # Convert column letter to column number
    return ord(letter.upper()) - 64  # Assuming A is 1, B is 2, etc.

def find_and_write_by_name(file_path, sheet_name, search_column, search_value, write_column, write_value):
    # Load the Excel workbook
    workbook = load_workbook(file_path)
    
    # Select the specified sheet
    sheet = workbook[sheet_name]
    
    # Determine the column index for search_column
    search_column_index = ord(search_column.upper()) - ord('A')
    
    # Calculate the target write column index
    if write_column >= 0:
        target_column_index = search_column_index + write_column + 1
    else:
        target_column_index = search_column_index + write_column
    
    # Adjust target_column_index if it is negative
    if target_column_index < 0:
        target_column_index = 0
    
    # Iterate over the rows in the sheet
    for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, min_col=search_column_index+1, max_col=search_column_index+1):
        for cell in row:
            if cell.value == search_value:
                # Write write_value to the corresponding column
                sheet.cell(row=cell.row, column=target_column_index + 1, value=write_value)
    
    # Save the changes
    workbook.save(file_path)

def fetch_data(excel_path, sheet_name, column_b_value, column_a_value, column_number):
    # Load the workbook
    workbook = load_workbook(excel_path)
    
    # Get the specified sheet
    sheet = workbook[sheet_name]
    
    # Iterate over rows in column B to find the matching value
    for row in range(1, sheet.max_row + 1):
        if sheet.cell(row=row, column=2).value == column_b_value:
            # Check the corresponding value in column A
            if sheet.cell(row=row, column=1).value == column_a_value:
                # Fetch data from the specified column
                data = sheet.cell(row=row, column=column_number).value
                return data
    
    # If the value is not found, return None
    return None

def fetch_data_arrays(excel_path, sheet_name, column_b_value, column_a_value, column_numbers):
    try:
        # Load the workbook
        workbook = load_workbook(excel_path, read_only=True)
        
        # Get the specified sheet
        sheet = workbook[sheet_name]
        
        # Initialize an empty list to hold the results
        results = []
        
        # Iterate over rows in column B to find the matching value
        for row in range(1, sheet.max_row + 1):
            if sheet.cell(row=row, column=2).value == column_b_value:
                # Check the corresponding value in column A
                if sheet.cell(row=row, column=1).value == column_a_value:
                    # Fetch data from the specified columns and add each value to the results
                    for col_num in column_numbers:
                        data = sheet.cell(row=row, column=col_num).value
                        results.append(data)
        
        # If no matching row is found, return a message indicating that
        if not results:
            return f"No matching row found for column B value '{column_b_value}' and column A value '{column_a_value}'."
        
        return results
    
    except Exception as e:
        # Handle any errors that might occur (e.g., file not found, sheet not found, etc.)
        return f"Error: {str(e)}"

def read_excel_cells_by_array(file_path, sheet_name, column_to_search, search_words, offsets):
    # Ensure both arrays are of the same length
    if len(search_words) != len(offsets):
        raise ValueError("The 'search_words' and 'offsets' arrays must have the same length.")
    
    # Load the workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    
    # Select the worksheet
    ws = wb[sheet_name]
    
    # Convert column letter to column index
    column_index = utils.column_index_from_string(column_to_search)
    
    # Array to store the results
    results = []
    
    # Iterate over search_words and corresponding offsets
    for search_word, offset in zip(search_words, offsets):
        # Iterate over the rows in the specified column
        found_value = None
        for row in ws.iter_rows(min_col=column_index, max_col=column_index, values_only=False):
            cell_value = row[0].value  # Access the first cell in the row
            
            # Check if the cell value matches the search word
            if cell_value == search_word:
                # Get the row number of the matching word
                row_number = row[0].row
                
                # Calculate the column index based on the offset
                col_number = column_index + int(offset)
                
                # Fetch the value from the corresponding column in the same row
                found_value = ws.cell(row=row_number, column=col_number).value
                break  # Exit the loop once the word is found
        
        # Append the found value (or None if not found) to the results array
        results.append(found_value)
    
    # Close the workbook
    wb.close()
    
    # Return the array of results
    return results

def read_excel_cells_default_offsets(file_path, sheet_name, column_to_search, search_words):
    # Load the workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    
    # Select the worksheet
    ws = wb[sheet_name]
    
    # Convert column letter to column index
    column_index = utils.column_index_from_string(column_to_search)
    
    # Array to store the results
    results = []
    
    # Define the default offsets (-1 and -3)
    default_offsets = [-1, -3]
    
    # Iterate over search_words
    for search_word in search_words:
        # Iterate over the rows in the specified column
        for row in ws.iter_rows(min_col=column_index, max_col=column_index, values_only=False):
            cell_value = row[0].value  # Access the first cell in the row
            
            # Check if the cell value matches the search word
            if cell_value == search_word:
                # Get the row number of the matching word
                row_number = row[0].row
                
                # Fetch values from the corresponding columns with offsets -1 and -3
                for offset in default_offsets:
                    col_number = column_index + offset
                    found_value = ws.cell(row=row_number, column=col_number).value
                    
                    # Append the found value (or None if not found) to the results array
                    results.append(found_value)
                
                # Move to the next search word once found
                break
    
    # Close the workbook
    wb.close()
    
    # Return the 1-D array of results
    return results

def find_section_range_arrays(excel_path, sheet_name, words):
    try:
        # Load the Excel workbook
        workbook = load_workbook(excel_path, read_only=True)
        
        # Select the worksheet by name
        worksheet = workbook[sheet_name]
        
        # Initialize an empty list to hold the results
        results = []
        
        # Iterate over the words and their corresponding column numbers
        for word in words:
            found = False
            # Iterate over rows in the worksheet
            for row in worksheet.iter_rows(values_only=True):
                if row[1] == word: 
                    # Assuming the section name is in column 0, start slide in column 1, and end slide in column 2
                    start_slide = row[2]  # Start slide number (column 1)
                    end_slide = row[3]    # End slide number (column 2)
                    results.append((start_slide, end_slide))  # Append the range for the section
                    found = True
                    break
            if not found:
                results.append((None, None))  # Append None if the section wasn't found
        
        return results
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        return []

def show_slide_ranges_from_sections(ppt_file, excel_path, sheet_name, section_ids):
    # Open the PowerPoint presentation
    presentation = Presentation(ppt_file)
    
    # Fetch start and end slide numbers for each section from the Excel sheet
    slide_ranges = find_section_range_arrays(excel_path, sheet_name, section_ids)
    
    # Validate and iterate over the slide ranges
    for i, (start_slide, end_slide) in enumerate(slide_ranges):
        if start_slide is None or end_slide is None:
            print(f"Section '{find_slide_nums_arrays_v2(excel_path, sheet_name, [section_ids[i]], [-1])}' not found in the Excel sheet.")
            continue
        
        # Adjust for 0-based indexing
        start_index = start_slide - 1
        end_index = end_slide - 1
        
        # Validate slide range
        if 0 <= start_index <= end_index < len(presentation.slides):
            # Iterate over the specified range and show each slide
            for slide_index in range(start_index, end_index + 1):
                slide = presentation.slides[slide_index]
                slide._element.set('show', '1')  # Sets the slide to be visible
        else:
            print(f"Invalid slide range for section '{section_ids[i]}': {start_slide} to {end_slide}.")
    
    # Save the modified presentation
    presentation.save(ppt_file)

def hide_slide_ranges_from_sections(ppt_file, excel_path, sheet_name, section_ids):
    # Open the PowerPoint presentation
    presentation = Presentation(ppt_file)
    
    # Fetch start and end slide numbers for each section from the Excel sheet
    slide_ranges = find_section_range_arrays(excel_path, sheet_name, section_ids)
    
    # Validate and iterate over the slide ranges
    for i, (start_slide, end_slide) in enumerate(slide_ranges):
        if start_slide is None or end_slide is None:
            print(f"Section '{find_slide_nums_arrays_v2(excel_path, sheet_name, [section_ids[i]], [-1])}' not found in the Excel sheet.")
            continue
        
        # Adjust for 0-based indexing
        start_index = start_slide - 1
        end_index = end_slide - 1
        
        # Validate slide range
        if 0 <= start_index <= end_index < len(presentation.slides):
            # Iterate over the specified range and show each slide
            for slide_index in range(start_index, end_index + 1):
                slide = presentation.slides[slide_index]
                slide._element.set('show', '0')  # Sets the slide to be visible
        else:
            print(f"Invalid slide range for section '{section_ids[i]}': {start_slide} to {end_slide}.")
    
    # Save the modified presentation
    presentation.save(ppt_file)

def open_presentation_relative_path(rp):
    absolute_path = relative_path(rp)
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(absolute_path)
    return presentation

def find_slide_index_by_title(presentation, title, start_index=1, direction="down"):
    # Get the total number of slides in the presentation
    num_slides = presentation.Slides.Count
    
    # Validate start_index
    if start_index < 1 or start_index > num_slides:
        raise ValueError("start_index is out of range")

    # Determine the range of iteration based on the direction
    if direction == "down":
        slide_indices = range(start_index, num_slides + 1)
    elif direction == "up":
        slide_indices = range(start_index, 0, -1)
    else:
        raise ValueError("direction must be 'up' or 'down'")
    
    # Iterate through slides based on the specified direction
    for i in slide_indices:
        slide = presentation.Slides(i)
        
        for shape in slide.Shapes:
            if shape.HasTextFrame:
                text_frame = shape.TextFrame
                if text_frame.HasText:
                    if title in text_frame.TextRange.Text:
                        return i
    return -1  # Return -1 if the title is not found in any slide

def find_slide_index_by_label(presentation, label, start_index=1):
    # Get the total number of slides in the presentation
    num_slides = presentation.Slides.Count

    # Iterate through each slide starting from the specified start_index
    if start_index < 1 or start_index > num_slides:
        raise ValueError("start_index is out of range")

    for i in range(start_index, num_slides + 1):
        slide = presentation.Slides(i)
        for shape in slide.Shapes:
            if shape.HasTextFrame:
                text_frame = shape.TextFrame
                if text_frame.HasText:
                    # Check if the label of the shape matches the provided label
                    if shape.Name == label:
                        return i
    return -1

def find_slide_indices_by_ordered_labels(presentation, labels, start_index=1):
    # Get the total number of slides in the presentation
    num_slides = presentation.Slides.Count

    # Initialize a dictionary to store label-to-slide index mapping
    label_to_slide_index = {label: -1 for label in labels}
    current_slide_index = start_index

    # Iterate through each label in sequence
    for label in labels:
        # Search for the label starting from the last found slide index
        
        for i in range(current_slide_index, num_slides + 1):
            slide = presentation.Slides(i)
            label_found = False

            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    if shape.Name == label:
                        # Update the mapping and set current_slide_index to next slide
                        label_to_slide_index[label] = i
                        current_slide_index = i + 1  # Move to next slide for next label
                        label_found = True
                        break

            if label_found:
                break  # Proceed to the next label

    # Return an array of slide indices in the order of the input labels
    return [label_to_slide_index[label] for label in labels]

def get_slide_ids_by_numbers(pptx_path, slide_numbers):
    # Load the presentation
    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)
    
    # Initialize the list to hold SlideIDs
    slide_ids = []

    for slide_num in slide_numbers:
        # Check if the slide number is within the valid range
        if slide_num < 1 or slide_num > num_slides:
            print(f"Error: Slide number {slide_num} is out of range.")
            slide_ids.append(None)  # Append None for invalid slide numbers
            continue
        
        # Get the SlideID for the slide number
        slide_id = prs.slides[slide_num - 1].slide_id  # Convert to 0-based index

        # Add the SlideID to the list
        slide_ids.append(slide_id)
    
    return slide_ids

def get_slide_ids_by_number_pairs(pptx_path, slide_number_pairs):
    # Load the presentation
    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)
    
    # Initialize the list to hold pairs of SlideIDs
    slide_id_pairs = []

    for slide_num1, slide_num2 in slide_number_pairs:
        # Check if both slide numbers are within the valid range
        if slide_num1 < 1 or slide_num1 > num_slides or slide_num2 < 1 or slide_num2 > num_slides:
            print(f"Error: Slide numbers ({slide_num1}, {slide_num2}) out of range.")
            slide_id_pairs.append((None, None))  # Append None for invalid pairs
            continue
        
        # Get the SlideID for each slide number
        slide_id1 = prs.slides[slide_num1 - 1].slide_id  # Convert to 0-based index
        slide_id2 = prs.slides[slide_num2 - 1].slide_id

        # Add the SlideID pair to the list
        slide_id_pairs.append((slide_id1, slide_id2))
    
    return slide_id_pairs

def get_slide_number_by_id(presentation, slide_id):
    # Iterate through slides to match Slide ID
    for slide in presentation.Slides:
        if slide.SlideID == slide_id:
            slide_number = slide.SlideIndex  # SlideIndex is the slide's position number
            break    
    return slide_number

def insert_image_to_slides_same_file(original_presentation_path, image_path):
    from pptx.util import Inches
    from PIL import Image
    
    prs = Presentation(original_presentation_path)
    image = Image.open(image_path)
    dpi = image.info.get("dpi", (72, 72))[0]
    width_pixels, height_pixels = image.size
    width_inches = width_pixels/dpi
    height_inches = height_pixels/dpi

    left = Inches(0)
    top = Inches(1.980315)
    width = Inches(width_inches)  
    height = Inches(height_inches)  

    for slide in prs.slides:
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        slide.shapes._spTree.insert(2, pic._element)
    prs.save(original_presentation_path)

def hide_slides(presentation, hide_array):
    for start_slide, end_slide in hide_array:
        # print(start_slide, end_slide)
        for i in range(start_slide, end_slide + 1):
            presentation.Slides(i).SlideShowTransition.Hidden = True

def show_slides(presentation, show_array):
    for start_slide, end_slide in show_array:
        # print(start_slide, end_slide)
        for i in range(start_slide, end_slide + 1):
            presentation.Slides(i).SlideShowTransition.Hidden = False

def find_Readings_Date (month, day):
    special_days = {
        1: [1, 2, 8, 16, 17, 18, 19, 21, 26],
        2: [12, 14, 22, 27],
        3: [8, 9, 12, 15, 17, 22, 24, 25, 27, 28, 29],
        4: [22, 23, 28, 29, 30],
        5: [1, 3, 4, 6, 10, 11, 12, 13, 22, 26, 30],
        6: [2],
        7: [13, 29],
        8: [23, 27, 30],
        9: [1, 10, 20, 24, 26],
        10: [2, 16, 30],
        11: [3, 5, 20],
        12: [3, 13, 17, 25, 26, 28, 29, 30],
        13: [1, 2, 3, 4, 6]  # Assuming 13th month special days
    }
    if month in special_days and day in special_days[month]:
        return month, day
    
    #28 هاتور
    hator28 = {
        1: [9],
        2: [23],
        3: [4, 14, 23],
        4: [10, 19, 27],
        5: [18, 23],
        6: [22, 24],
        7: [2, 5, 17],
        8: [19],
        9: [11, 17, 27],
        10: [3, 14],
        11: [10, 13],
        12: [11, 24],
        13: [5]
    }
    if month in hator28 and day in hator28[month]:
        return 3, 28
    
    # 30 طوبة
    toba30 = {
        1: [5, 10, 20, 29],
        2: [1, 6],
        3: [10],
        4: [26],
        5: [24, 29],
        7: [15, 26],
        8: [6, 8, 11, 25],
        11: [1, 17, 28],
        12: [2, 6, 15, 21]
    }
    if month in toba30 and day in toba30[month]:
        return 5, 30
    
    # 3 أبيب
    abib3 = {
        1: [3, 7],
        2: [2, 11, 17],
        3: [1, 13, 30],
        4: [6],
        5: [2, 8],
        6: [7],
        7: [22],
        8: [3, 12],
        9: [7, 30],
        10: [28],
        12: [18]
    }
    if month in abib3 and day in abib3[month]:
        return 11, 3
    
    # 8 توت
    tot8 = {
        1: [4, 6, 25],
        2: [21],
        4: [5, 20],
        5: [15],
        6: [26],
        7: [23],
        8: [5, 7],
        9: [5],
        10: [9, 20, 26],
        11: [26],
        12: [4, 22]
    }
    if month in tot8 and day in tot8[month]:
        return 1, 8
    
    # 17 هاتور
    hator17 = {
        1: [13, 30],
        2: [3, 5, 9],
        3: [6, 21],
        4: [1, 15],
        5: [7],
        6: [14, 18, 27],
        9: [12, 28],
        12: [11, 24],
    }
    if month in hator17 and day in hator17[month]:
        return 3, 17
    
    # 29 هاتور
    hator29 = {
        2: [16, 18],
        3: [2],
        6: [11, 20],
        7: [3, 12, 16, 20, 24],
        8: [14, 22],
        9: [4],
        10: [18],
    }
    if month in hator29 and day in hator29[month]:
        return 3, 29

    # 22 هاتور
    hator22 = {
        1: [27, 28],
        3: [26],
        5: [19],
        7: [1, 7],
        8: [4, 13],
        10: [10, 22],
        11: [19, 30],
        12: [20]
    }
    if month in hator22 and day in hator22[month]:
        return 3, 22
    
    # 20 بشنس
    bashans20 = {
         2: [7, 13, 25],
         4: [12],
         5: [14, 25],
         6: [3, 12, 19],
         7: [9],
         8: [1, 10],
         10 : [29]
    }
    if month in bashans20 and day in bashans20[month]:
        return 9, 20

    # 27 برمودة
    barmoda27 = {
         1: [11, 22],
         4: [14], 
         5: [16, 27],
         6: [13, 17, 23],
         8: [18],
         10: [11],
         11: [11, 22]
    }
    if month in barmoda27 and day in barmoda27[month]:
        return 8, 27
    
    # 1 طوبة
    toba1 = {
        1: [15, 24],
        2: [26],
        5: [20],
        7: [19],
        8: [15, 29],
        9: [3, 22],
        10: [1, 25],
        11: [6]
    }
    if month in toba1 and day in toba1[month]:
        return 5, 1
    
    # 20 ابيب
    abib20 = {
        2: [15, 19],
        6: [25, 28],
        7: [18],
        8: [2, 20, 24, 26],
        10: [6]
    }
    if month in abib20 and day in abib20[month]:
        return 11, 20
    
    # 3 مسرى
    mesra3 = {
        1: [14],
        4: [17],
        6: [9],
        9: [2, 14, 21, 29],
        10: [24],
        11: [15]
    }
    if month in mesra3 and day in mesra3[month]:
        return 11, 3
    
    # 15 هاتور
    hator15 = {
        5: [28],
        6: [6],
        10: [4, 15],
        11: [4, 12, 24],
        12: [1, 9]
    }
    if month in hator15 and day in hator15[month]:
        return 3, 15
    
    # 5 ابيب 
    abib5 = {
        3: [18],
        4: [4, 21],
        6: [10, 21],
        7: [8],
        8: [17],
        11: [18, 29]
    }
    if month in abib5 and day in abib5[month]:
        return 11, 5
    
    # 22 طوبة
    toba22 = {
        2: [24],
        4: [9],
        8: [9],
        9: [13, 19],
        10: [5, 17, 23]
    }
    if month in toba22 and day in toba22[month]:
        return 5, 22
    
    # 23 برمودة
    baramoda23 = {
        3: [7],
        9: [6, 8],
        10: [7, 19],
        11: [27],
        12: [23, 27]
    }
    if month in baramoda23 and day in baramoda23[month]:
        return 8, 23
    
    # 27 بابة
    baba27 = {
        5: [9],
        6: [5],
        7: [27],
        9: [18],
        11: [8, 14],
        12: [19]
    }
    if month in baba27 and day in baba27[month]:
        return 2, 27

    # 25 هاتور
    hator25 = {
        2: [8, 28],
        5: [5],
        8: [28],
        9: [25],
        11: [25]
    }
    if month in hator25 and day in hator25[month]:
        return 3, 25

    # 27 هاتور
    hator27 = {
        2: [29],
        3: [5],
        4: [16],
        7: [6],
        11: [23],
        12: [10]
    }
    if month in hator27 and day in hator27[month]:
        return 3, 27

    # 30 برمودة
    baramoda30 = {
        2: [30],
        3: [20],
        9: [23],
        10: [27],
        11: [2, 9]
    }
    if month in baramoda30 and day in baramoda30[month]:
        return 8, 30
    
    # 13 برمهات 
    barmhat13 = {
        2: [4, 10],
        3: [19],
        4: [8],
        11: [21]  
    }
    if month in barmhat13 and day in barmhat13[month]:
        return 7, 13

    # 16 بؤونة
    bo2ona16 = {
        3: [16],
        4: [2, 11],
        5: [17],
        6: [15]
    }
    if month in bo2ona16 and day in bo2ona16[month]:
        return 10, 16
    
    # 4 طوبة
    toba4 = {
        4: [24],
        6: [29],
        8: [16],
        9: [16],
        11: [16]
    }
    if month in toba4 and day in toba4[month]:
        return 5, 4
    
    # 9 هاتور
    hator9 = {
        1: [12],
        3: [3],
        6: [1],
        7: [4]
    }
    if month in hator9 and day in hator9[month]:
        return 3, 9
    
    # 21 توت
    tot21 = {
        1: [23],
        7: [11, 14],
        8: [21]
    }
    if month in tot21 and day in tot21[month]:
        return 1, 21

    # 2 امشير
    amshyr2 = {
        2: [20],
        4: [7, 25],
        11: [7]
    }
    if month in amshyr2 and day in amshyr2[month]:
        return 6, 2

    # 1 بشنس
    bashans1 = {
        4: [3],
        5: [21],
        10: [21],
        12: [16]
    }
    if month in bashans1 and day in bashans1[month]:
        return 9, 1 
    
    # 26 توت
    tot26 = {
        3: [11],
        4: [13],
        6: [16],
        12: [7]
    }
    if month in tot26 and day in tot26[month]:
        return 1, 26
    
    # 16 توت
    tot16 = {
        7: [28],
        9: [9]
    }
    if month in tot16 and day in tot16[month]:
        return 1, 16
    
    # 17 توت
    tot17 = {
        7: [10],
        12: [12]
    }
    if month in tot17 and day in tot17[month]:
        return 1, 17

    # 22 بابة
    baba22 = {
        4: [18],
        6:[4]
    }
    if month in baba22 and day in baba22[month]:
        return 2, 22
    
    # 12 بابة
    baba12 = {
        9: [15]
    }
    if month in baba12 and day in baba12[month]:
        return 2, 12

    # 14 بابة
    baba14 = {
        7: [25]
    }
    if month in baba14 and day in baba14[month]:
        return 2, 14

    # 12 هاتور
    hator12 = {
        10: [12]
    }
    if month in hator12 and day in hator12[month]:
        return 3, 12
    
    # 30 برمهات
    kiahk22 = {
        7: [30]
    }
    if month in kiahk22 and day in kiahk22[month]:
        return 4, 22

    # 6 طوبة
    toba6 = {
        6: [8]
    }
    if month in toba6 and day in toba6[month]:
        return 5, 6

    # 26 طوبة
    toba26 = {
        12: [5]
    }
    if month in toba26 and day in toba26[month]:
        return 5, 26

    # 2 بؤونة
    bo2ona2 = {
        6: [30]
    }
    if month in bo2ona2 and day in bo2ona2[month]:
        return 10, 2

    # 13 مسرى
    mesra13 = {
        7: [21]
    }
    if month in mesra13 and day in mesra13[month]:
        return 12, 13

    # 10 بشنس
    bashans10 = {
        12: [8]
    }
    if month in bashans10 and day in bashans10[month]:
        return 9, 10
    
    # 24 بشنس
    bashans24 = {
        10: [8]
    }
    if month in bashans24 and day in bashans24[month]:
        return 9, 24

    # 3 نسئ
    nsy3 = {
        10: [13]
    }
    if month in nsy3 and day in nsy3[month]:
        return 13, 3

def get_slide_ids_by_number(pptx_path, slide_num):
    # Load the presentation
    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)
    
    if slide_num < 1 or slide_num > num_slides:
        print(f"Error: Slide number {slide_num} is out of range.")
        return
        
    # Get the SlideID for the slide number
    slide_id = prs.slides[slide_num - 1].slide_id  # Convert to 0-based index    
    return slide_id

def move_sections(presentation, move_section_names, target_section_names):
    for move_section_name, target_section_name in zip(move_section_names, target_section_names):
        sections = {presentation.SectionProperties.Name(i): i for i in range(1, presentation.SectionProperties.Count + 1)}
        move_index = sections[move_section_name]
        target_index = sections[target_section_name]
        if move_index < target_index:
            target_index -= 1
        presentation.SectionProperties.Move(move_index, target_index + 1)

def move_sections_v2(presentation, move_section_ids, target_section_ids):
    for move_section_id, target_section_id in zip(move_section_ids, target_section_ids):
        sections = {presentation.SectionProperties.SectionID(i): i for i in range(1, presentation.SectionProperties.Count + 1)}
        move_index = sections[move_section_id]
        target_index = sections[target_section_id]
        if move_index < target_index:
            target_index -= 1
        presentation.SectionProperties.Move(move_index, target_index + 1)

def move_sections_range(presentation, start_section_name, end_section_name, target_section_name):
    sections = {presentation.SectionProperties.Name(i): i for i in range(1, presentation.SectionProperties.Count + 1)}
    if start_section_name not in sections or end_section_name not in sections or target_section_name not in sections:
        raise ValueError("One or more specified section names are invalid.")
    start_index = sections[start_section_name]
    end_index = sections[end_section_name]
    target_index = sections[target_section_name]
    if start_index > end_index:
        raise ValueError("Start section must come before or be the same as the end section.")
    move_indices = list(range(start_index, end_index + 1))
    if target_index >= start_index and target_index <= end_index:
        raise ValueError("Target section cannot be within the range of sections to move.")
    if target_index > end_index:
        target_index -= len(move_indices)
    for move_index in move_indices:
        presentation.SectionProperties.Move(move_index, target_index + 1)
        target_index += 1

def run_vba_with_slide_id_bakr_aashya(excel, sheet, prs, presentation):

    slide = find_slide_num_v2(excel, sheet, '{A5B9CE2F-90E3-44D7-B22F-CAE6783C8E2F}', 2, 1)

    slide_id = get_slide_ids_by_number(prs, slide)

    # Access the VBA project
    vba_project = presentation.VBProject
    modules = vba_project.VBComponents

    # Add a new module to the VBA project
    new_module = modules.Add(1)  # 1 corresponds to a standard module

    vba_code = f"""
Dim visitedSlides As Collection ' Global collection to track visited slides

Sub OnSlideShowPageChange()
    Dim currentSlideNumber As Integer
    Dim targetShape As Shape

    ' Initialize the visitedSlides collection if it hasn't been created yet
    If visitedSlides Is Nothing Then
        Set visitedSlides = New Collection
    End If

    ' Get the current slide number in the slideshow view
    currentSlideID = ActivePresentation.SlideShowWindow.View.Slide.SlideID

    ' Check if the slide has already been visited (i.e., hyperlink followed)
    On Error Resume Next
    visitedSlides.Item currentSlideID
    If Err.Number = 0 Then
        ' Slide has already been visited; exit without doing anything
        Exit Sub
    End If
    On Error GoTo 0

    ' Use Select Case to handle actions on specific slides by slide number
    Select Case currentSlideID
        Case {slide_id} ' Replace with the slide number for the first target slide
            ' Attempt to locate and "click" the target shape
            Set targetShape = GetShapeByName("TextBox 2") ' Replace with your shape name
            
            If Not targetShape Is Nothing Then
                If targetShape.ActionSettings(ppMouseClick).Action = ppActionHyperlink Then
                    ' Follow the hyperlink
                    targetShape.ActionSettings(ppMouseClick).Hyperlink.Follow
                End If
            End If
    End Select
End Sub

' Helper function to get a shape by name on the current slide
Function GetShapeByName(shapeName As String) As Shape
    On Error Resume Next
    Set GetShapeByName = ActivePresentation.SlideShowWindow.View.Slide.Shapes(shapeName)
    On Error GoTo 0
End Function

    """
    # Add the generated code to the new module
    new_module.CodeModule.AddFromString(vba_code)

    # Set up the slideshow to call OnSlideShowPageChange on each slide change
    presentation.SlideShowSettings.Run()

    # Optionally run the macro immediately to initialize
    presentation.Application.Run("OnSlideShowPageChange")

    presentation.SlideShowWindow.View.Exit()

def replacefile(old_file, new_file):
    from shutil import copy2
    from os import path, remove
    try:
        # If the old file exists, delete it
        if path.exists(old_file):
            remove(old_file)
        
        # Copy the new file to the location of the old file
        copy2(new_file, old_file)

    except Exception as e:
        # Print any errors that occur during the deletion and copying process
        print(f"Error: {str(e)}")

def elzoksologyat (excel_path, season, bakerOR3ashyaORtasbha):
    replacefile(relative_path(r"الذكصولوجيات.pptx"), relative_path(r"Data\CopyData\الذكصولوجيات.pptx"))
    pptx_file = relative_path(r"الذكصولوجيات.pptx")
    sheet = "الذكصولوجيات"
    if bakerOR3ashyaORtasbha == "باكر":
        bakerOR3ashyaORtasbha = '{9621F9CE-ABC8-4FF6-A8C6-3AA9D24690A0}'
    elif bakerOR3ashyaORtasbha == "عشية":
        bakerOR3ashyaORtasbha = '{267B00F5-E8C1-4DF6-A5CB-DFF5531064E8}'
    else:
        bakerOR3ashyaORtasbha = '{38C049BC-0822-439F-B5A3-C6094A6A24B1}'
    match(season):
        case 1: show_slide_ranges_from_sections(pptx_file, excel_path, sheet, ["{9A902651-94A5-4D6A-83A0-BF404F380CD5}", bakerOR3ashyaORtasbha])
        case 2: show_slide_ranges_from_sections(pptx_file, excel_path, sheet, ["{264E4307-BE55-4799-9151-3D149372B553}", bakerOR3ashyaORtasbha])
        case 29: show_slide_ranges_from_sections(pptx_file, excel_path, sheet, ["{DC61EC21-9EF0-4E7C-8E4E-CD4269024AE6}", bakerOR3ashyaORtasbha])
        case 5: show_slide_ranges_from_sections(pptx_file, excel_path, sheet, ["كيهك 1", "كيهك 2", "كيهك 3", "كيهك 4", "كيهك 5", "كيهك 6"])
        case 15 | 15.1 | 15.2 | 15.3 | 15.4 | 15.5 | 15.6 | 15.7 | 15.8 | 15.9 | 15.11: show_slide_ranges_from_sections(pptx_file, excel_path, sheet, ["الصوم الكبير 1", "الصوم الكبير 2", "الصوم الكبير 3", "الصوم الكبير 4", "الصوم الكبير 5"])
        case default: show_slide_ranges_from_sections(pptx_file, excel_path, sheet, [bakerOR3ashyaORtasbha])

def find_section_Ids_with_names(excel_path, sheet, names):
    try:
        # Load the Excel file
        wb = load_workbook(excel_path, data_only=True)
        ws = wb[sheet]
        section_ids = []
        not_found_names = []
        for name in names:
            found = False
            for row in reversed(list(ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=4, values_only=True))):
                if row[0] == name:
                    section_ids.append(row[1])
                    found = True
                    break
            if not found:
                not_found_names.append(name)
        if not_found_names:
            section_ids.extend(not_found_names)
        return section_ids
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        return []

def find_section_names_with_ids(excel, sheet, ids):
    try:
        # Load the Excel file
        wb = load_workbook(excel, data_only=True)
        ws = wb[sheet]
        section_names = []
        for row in reversed(list(ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=4, values_only=True))):
            for id in ids:
                if row[1] == id:
                    section_names.append(row[0])
        return section_names
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        return []

def show_hide_insertImage_replaceText(ppt_file, excel_path, sheet_name, 
                               show_sections=None, hide_sections=None, 
                               show_sections_ranges=None, hide_sections_ranges=None,
                               image_path=None, new_Text=None):
    """
    Combined function to show/hide slides and optionally insert an image.
    
    Parameters:
    - ppt_file: Path to PowerPoint file
    - excel_path: Path to Excel file with section data
    - sheet_name: Sheet name in Excel file
    - show_sections: 1D array of section IDs to show
    - hide_sections: 1D array of section IDs to hide
    - show_sections_ranges: 2D array of section ID ranges to show
    - hide_sections_ranges: 2D array of section ID ranges to hide
    - image_path: Optional path to image to insert
    - image_position: Optional tuple of (left, top, width, height) in inches for image positioning
                    If not provided and image_path is given, default position is used
    """
    
    from pptx.util import Inches
    from PIL import Image
    from pptx import Presentation
    
    # Open the PowerPoint presentation
    presentation = Presentation(ppt_file)

    # Function to process both 1D and 2D arrays into a set of slide indices
    def process_sections(section_data):
        slides_set = set()
        if not section_data:
            return slides_set
            
        if isinstance(section_data[0], list):  # 2D array case (ranges)
            for section_range in section_data:
                if len(section_range) >= 2:
                    slide_ranges = find_section_range_arrays(excel_path, sheet_name, section_range)
                    if slide_ranges:
                        start_slide = slide_ranges[0][0]  # Start of the first section
                        end_slide = slide_ranges[-1][1]  # End of the last section
                        if start_slide and end_slide:
                            slides_set.update(range(start_slide - 1, end_slide))
        else:  # 1D array case (individual sections)
            slide_ranges = find_section_range_arrays(excel_path, sheet_name, section_data)
            slides_set.update(
                slide for start, end in slide_ranges if start and end for slide in range(start - 1, end)
            )
        return slides_set

    # Process show and hide sections (1D and 2D separately)
    show_slides_set = process_sections(show_sections) if show_sections else set()
    hide_slides_set = process_sections(hide_sections) if hide_sections else set()
    show_ranges_set = process_sections(show_sections_ranges) if show_sections_ranges else set()
    hide_ranges_set = process_sections(hide_sections_ranges) if hide_sections_ranges else set()

    # Combine sets to get final show and hide lists
    final_show_slides = show_slides_set.union(show_ranges_set)
    final_hide_slides = hide_slides_set.union(hide_ranges_set)

    # Handle image insertion if image_path is provided
    if image_path:
        # Load image dimensions
        image = Image.open(image_path)
        
        image = Image.open(image_path)
        dpi = image.info.get("dpi", (72, 72))[0]
        width_pixels, height_pixels = image.size
        width_inches = width_pixels / dpi
        height_inches = height_pixels / dpi

        # Image position and size
        left = Inches(0)
        top = Inches(1.980315)
        width = Inches(width_inches)
        height = Inches(height_inches)

    if new_Text:
        # Text replacements
        find_replace_pairs = [
            ("لأنك أتيت", new_Text[0]),
            ("ak`i", new_Text[1]),
            ("آك إي", new_Text[2]),
        ]

    # Iterate over all slides and perform operations
    for i, slide in enumerate(presentation.slides):
        # Show or hide slide based on its index
        if i in final_show_slides:
            slide._element.set('show', '1')  # Show slide
        elif i in final_hide_slides:
            slide._element.set('show', '0')  # Hide slide

        # Replace text within the slide if new_Text is provided
        if new_Text:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for find, replace in find_replace_pairs:
                                if find in run.text:
                                    run.text = run.text.replace(find, replace)

        # Insert the image if path was provided
        if image_path:
            pic = slide.shapes.add_picture(image_path, left, top, width, height)
            slide.shapes._spTree.insert(2, pic._element)

    # Save the modified presentation
    presentation.save(ppt_file)    

# excel = relative_path(r"بيانات القداسات.xlsx")
# sheet = "رفع بخور"
# arr = ["تكملة في حضور الاسقف", "طوبه هينا الكبيرة", "مارو اتشاسف"]
# print(find_section_Ids_with_names(excel, sheet, arr))
# arr2 = ['{2BCF4F8C-25F0-43C5-B224-6528B2EA3F2F}', '{F76B0D75-0474-45B5-B79F-7416F354543A}',
#         '{E2968C91-5339-499C-9812-DECCCF58A2CD}', '{62A12AF8-CB6D-4CC5-9DB0-B73A7C24E2AD}', 
#         '{B74DBB8C-2B2D-46E4-9508-DA46008D19A4}', '{A9183893-7B7E-459F-8547-F7A8F7D2D521}', 
#         '{670DAA94-A6C9-4CCD-B4E2-958C71CD3E44}']
# print(find_section_names_with_ids(excel, sheet, arr2))
# print(find_slide_num_v2(excel, sheet, "تكملة للملاك ميخائيل 1", 1, 1))
# print(find_slide_nums_arrays_v2(excel, sheet, ['{BC7E3DCD-6AA8-44CC-B8AF-BC3E2BC71B5A}', '{BC7E3DCD-6AA8-44CC-B8AF-BC3E2BC71B5A}', '{A20DA654-32F7-4B4C-96CB-C76232EB96E8}', '{A20DA654-32F7-4B4C-96CB-C76232EB96E8}'], 2, [1, 2, 1, 2]))

# pptx_file = r"قداس.pptx"  # Path to your PowerPoint file
# image_file = r"Data\Designs\القيامة.png"  # Path to the image file
# insert_image_to_slides_same_file(pptx_file, image_file)

# from datetime import datetime
# start = datetime.now()
# ppt_file=relative_path(r"قداس.pptx")
# excel_path=relative_path(r"بيانات القداسات.xlsx")
# sheet_name="القداس"
# show_section_ids=['{BBBAC16F-044D-4F33-8068-620F498B59CD}', '{072F3D96-A6C8-405F-9A23-7CCA1B2F13FF}', '{670DAA94-A6C9-4CCD-B4E2-958C71CD3E44}', '{20F525FD-C708-4DDD-8E40-FE502EFEBDDE}', '{03E2AC57-01DD-4702-A7A7-186D0E009F55}', '{59DBF0F6-1D86-41E8-B37A-8AA2368AA8AB}', '{E6CBA825-E339-438B-84B4-326FC5C299C1}', '{8DD599A1-D7AC-4AA8-A52B-31BFD527E68E}', '{DEDC0CCA-3854-4E18-8CB2-5D6FEC5BABCC}', '{D95C2E5C-8772-445E-AE3E-2F50770CFC61}', '{B7D98377-B994-4654-B49C-DE10E0DDE4F1}', '{C2F28915-B86E-4596-8EB2-7455EF4E91BD}', '{42181297-997B-4C4C-B43B-4E9D8A23858D}']
# hide_section_ids=['{E107D25B-A642-458E-A4F3-B73FDB564A7C}', '{4D2B15D5-C978-467C-9D6C-726FE25128B8}', '{D5DB63D0-39EE-49CE-8855-58CE02719834}', '{31685B5B-48C4-437E-858C-CF8D225C0C26}', '{CD4B95FF-0E0E-42D3-8DC4-224C3DD732F7}', '{F13A48F2-238D-4617-B84E-9B0A694D9A18}', '{06D592C8-05BF-4B7C-86F7-FDAB3FAB5FB1}', '{681FF6A7-4230-4171-8F41-83FD64E8C960}', '{507EFD97-98F8-4376-848B-20D72E16D2C1}', '{B9A30F5E-0C89-471B-A99A-23DBE7F58504}']
# image_path=r"Data\Designs\الميلاد.png"
# newText = ["لأنك ولدت", "aumack", "اف ماسك"]

# # show_slide_ranges_from_sections(ppt_file, excel_path, sheet_name, show_section_ids)
# # hide_slide_ranges_from_sections(ppt_file, excel_path, sheet_name, hide_section_ids)
# # insert_image_to_slides_same_file(ppt_file, image_path)
# # milad(ppt_file)
# show_hide_insertImage_replaceText(ppt_file, excel_path, sheet_name, show_section_ids, hide_section_ids, image_path, newText)
# print(datetime.now() - start)
